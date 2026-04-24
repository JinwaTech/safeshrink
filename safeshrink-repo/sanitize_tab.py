"""
文档脱敏标签页
功能：检测并脱敏敏感信息
"""

from PySide6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QLabel,
    QPushButton, QTextEdit, QCheckBox, QFrame,
    QFileDialog, QMessageBox, QListWidget, QListWidgetItem,
    QGroupBox, QScrollArea, QRadioButton
)
from PySide6.QtCore import Qt
from pathlib import Path


class SanitizeTab(QWidget):
    def __init__(self):
        super().__init__()
        self.current_file = None
        self.original_content = None  # 原始内容（用于撤销）
        self.processed_content = None  # 处理后内容
        self.detected_items = []
        self._scene_ready = False  # 场景预设锁，防止信号级联触发
        
        # 原生文档对象（保持格式用）
        self._native_doc = None  # docx Document 对象
        self._native_docx_ext = None  # 原文件扩展名
        
        self.setup_ui()
        self._scene_ready = True
        self.scene_general.setChecked(True)
        self.load_settings()
    
    def load_settings(self):
        """从设置文件加载脱敏模式偏好（仅影响脱敏模式，不影响类型勾选）"""
        import json
        from pathlib import Path
        settings_file = Path(__file__).parent / 'settings.json'
        if settings_file.exists():
            try:
                with open(settings_file, 'r', encoding='utf-8') as f:
                    settings = json.load(f)
                mode = settings.get('sanitize_mode', 'mask')
                self.radio_mask.setChecked(mode == 'mask')
                self.radio_tags.setChecked(mode == 'tags')
            except:
                pass
        
    def setup_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 16, 0, 0)
        layout.setSpacing(16)
        
        # 文件选择（固定顶部）
        file_section = self.create_file_section()
        layout.addWidget(file_section)
        
        # 可滚动内容区
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setFrameShape(QFrame.Shape.NoFrame)
        
        scroll_content = QWidget()
        scroll_layout = QVBoxLayout(scroll_content)
        scroll_layout.setContentsMargins(0, 0, 0, 0)
        scroll_layout.setSpacing(16)
        
        # 敏感信息类型选择
        types_section = self.create_types_section()
        scroll_layout.addWidget(types_section)
        
        # 内容区域
        content_section = self.create_content_section()
        scroll_layout.addWidget(content_section, 1)
        
        scroll.setWidget(scroll_content)
        layout.addWidget(scroll, 1)
        
        # 操作按钮（固定底部）
        actions_section = self.create_actions_section()
        layout.addWidget(actions_section)
        
    def create_file_section(self):
        frame = QFrame()
        frame.setObjectName("card")
        layout = QHBoxLayout(frame)
        layout.setContentsMargins(20, 16, 20, 16)
        
        self.file_label = QLabel("未选择文件")
        self.file_label.setObjectName("subtitle")
        
        btn_browse = QPushButton("选择文件...")
        btn_browse.clicked.connect(self.browse_file)
        
        layout.addWidget(self.file_label)
        layout.addStretch()
        layout.addWidget(btn_browse)
        
        return frame
    
    def create_types_section(self):
        frame = QFrame()
        frame.setObjectName("card")
        layout = QVBoxLayout(frame)
        layout.setContentsMargins(20, 12, 20, 12)
        layout.setSpacing(8)
        
        title = QLabel("脱敏类型")
        title.setStyleSheet("font-weight: 600; font-size: 14px;")
        layout.addWidget(title)

        # 场景预设
        scene_box = QGroupBox("场景预设")
        scene_layout = QHBoxLayout(scene_box)
        scene_layout.setSpacing(16)
        self.scene_general = QRadioButton("通用文档")
        self.scene_gov = QRadioButton("党政公文")
        self.scene_finance = QRadioButton("金融合同")
        self.scene_medical = QRadioButton("医疗档案")
        self.scene_edu = QRadioButton("教育材料")
        self.scene_custom = QRadioButton("自定义")
        self.scene_general.toggled.connect(lambda: self._apply_scene('general'))
        self.scene_gov.toggled.connect(lambda: self._apply_scene('gov'))
        self.scene_finance.toggled.connect(lambda: self._apply_scene('finance'))
        self.scene_medical.toggled.connect(lambda: self._apply_scene('medical'))
        self.scene_edu.toggled.connect(lambda: self._apply_scene('edu'))
        self.scene_custom.toggled.connect(lambda: self._apply_scene('custom'))
        scene_layout.addWidget(self.scene_general)
        scene_layout.addWidget(self.scene_gov)
        scene_layout.addWidget(self.scene_finance)
        scene_layout.addWidget(self.scene_medical)
        scene_layout.addWidget(self.scene_edu)
        scene_layout.addWidget(self.scene_custom)
        scene_layout.addStretch()
        layout.addWidget(scene_box)

        # 脱敏模式选择
        mode_layout = QHBoxLayout()
        mode_label = QLabel("脱敏模式:")
        mode_label.setStyleSheet("color: #666; font-size: 12px;")
        
        self.radio_mask = QCheckBox("简单遮蔽 (***)")
        self.radio_mask.setChecked(True)
        self.radio_mask.setToolTip("将敏感信息替换为 ***")
        
        self.radio_tags = QCheckBox("语义标签 (<PERSON_NAME.1>)")
        self.radio_tags.setToolTip("用语义标签替代，大模型能理解类型但看不到内容")
        
        # 互斥选择
        self.radio_mask.stateChanged.connect(lambda: self.radio_tags.setChecked(not self.radio_mask.isChecked()))
        self.radio_tags.stateChanged.connect(lambda: self.radio_mask.setChecked(not self.radio_tags.isChecked()))
        
        mode_layout.addWidget(mode_label)
        mode_layout.addWidget(self.radio_mask)
        mode_layout.addWidget(self.radio_tags)
        mode_layout.addStretch()
        layout.addLayout(mode_layout)
        
        # 创建可滚动的内容容器
        scroll_content = QWidget()
        scroll_layout = QVBoxLayout(scroll_content)
        scroll_layout.setContentsMargins(0, 0, 0, 0)
        scroll_layout.setSpacing(8)
        
        # 使用水平布局分组
        types_layout = QHBoxLayout()
        
        # 个人敏感信息组
        personal_group = QVBoxLayout()
        group1_label = QLabel("个人敏感信息:")
        group1_label.setStyleSheet("color: #8b92a5; font-size: 12px;")
        
        self.chk_phone = QCheckBox("手机号 (138****8888)")
        self.chk_phone.setChecked(True)
        
        self.chk_email = QCheckBox("邮箱 (ab***@domain.com)")
        self.chk_email.setChecked(True)
        
        self.chk_idcard = QCheckBox("身份证号 (3301***********4)")
        self.chk_idcard.setChecked(True)
        
        self.chk_bank = QCheckBox("银行卡号 (1234****5678)")
        self.chk_bank.setChecked(True)
        
        self.chk_ip = QCheckBox("IP地址 (xxx.xxx.xxx.xxx)")
        self.chk_ip.setChecked(True)
        
        self.chk_passport = QCheckBox("护照号 (G********1)")
        self.chk_passport.setChecked(True)
        
        self.chk_mac = QCheckBox("Mac 地址")
        self.chk_mac.setChecked(True)
        
        self.chk_imei = QCheckBox("IMEI 设备号")
        self.chk_imei.setChecked(True)
        
        self.chk_plate = QCheckBox("车牌号 (浙A*****5)")
        self.chk_plate.setChecked(True)
        
        self.chk_social = QCheckBox("社保卡号")
        self.chk_social.setChecked(True)
        
        personal_group.addWidget(group1_label)
        personal_group.addWidget(self.chk_phone)
        personal_group.addWidget(self.chk_email)
        personal_group.addWidget(self.chk_idcard)
        personal_group.addWidget(self.chk_bank)
        personal_group.addWidget(self.chk_ip)
        personal_group.addWidget(self.chk_passport)
        personal_group.addWidget(self.chk_mac)
        personal_group.addWidget(self.chk_imei)
        personal_group.addWidget(self.chk_plate)
        personal_group.addWidget(self.chk_social)
        personal_group.addStretch()
        
        # 商业敏感信息组
        business_group = QVBoxLayout()
        group2_label = QLabel("商业敏感信息:")
        group2_label.setStyleSheet("color: #8b92a5; font-size: 12px;")
        
        self.chk_credit = QCheckBox("社会信用代码 (9133***XXXX)")
        self.chk_credit.setChecked(True)
        
        self.chk_license = QCheckBox("营业执照号")
        self.chk_license.setChecked(True)
        
        self.chk_contract = QCheckBox("合同编号 (HT****)")
        self.chk_contract.setChecked(True)
        
        self.chk_amount = QCheckBox("投标/成交价 (¥***)")
        self.chk_amount.setChecked(True)
        
        self.chk_phone2 = QCheckBox("固定电话")
        self.chk_phone2.setChecked(True)
        
        self.chk_account_permit = QCheckBox("开户许可证号")
        self.chk_account_permit.setChecked(True)
        
        self.chk_purchase_order = QCheckBox("采购/订单编号")
        self.chk_purchase_order.setChecked(True)
        
        self.chk_fax = QCheckBox("传真号")
        self.chk_fax.setChecked(True)
        
        self.chk_employee_id = QCheckBox("工号/学号")
        self.chk_employee_id.setChecked(True)
        
        self.chk_project_code = QCheckBox("项目代号")
        self.chk_project_code.setChecked(True)
        
        self.chk_postal = QCheckBox("邮编")
        self.chk_postal.setChecked(True)
        
        business_group.addWidget(group2_label)
        business_group.addWidget(self.chk_credit)
        business_group.addWidget(self.chk_license)
        business_group.addWidget(self.chk_contract)
        business_group.addWidget(self.chk_amount)
        business_group.addWidget(self.chk_phone2)
        business_group.addWidget(self.chk_account_permit)
        business_group.addWidget(self.chk_purchase_order)
        business_group.addWidget(self.chk_fax)
        business_group.addWidget(self.chk_employee_id)
        business_group.addWidget(self.chk_project_code)
        business_group.addWidget(self.chk_postal)
        business_group.addStretch()
        
        # 专用类型组
        special_group = QVBoxLayout()
        group3_label = QLabel("专用类型:")
        group3_label.setStyleSheet("color: #8b92a5; font-size: 12px;")
        
        gov_label = QLabel("  党政公文:")
        gov_label.setStyleSheet("font-size: 11px; color: #6b7280;")
        
        self.chk_docnum = QCheckBox("公文份号 (№******)")
        self.chk_docnum.setChecked(True)
        
        self.chk_doclevel = QCheckBox("密级标注 (绝密★***年)")
        self.chk_doclevel.setChecked(True)
        
        self.chk_docref = QCheckBox("公文文号 (〔2024〕*字第***号)")
        self.chk_docref.setChecked(True)
        
        med_label = QLabel("  医疗档案:")
        med_label.setStyleSheet("font-size: 11px; color: #6b7280;")
        
        self.chk_medicare = QCheckBox("医保卡号")
        self.chk_medicare.setChecked(True)
        
        self.chk_medical_record = QCheckBox("病历号/门诊号")
        self.chk_medical_record.setChecked(True)
        
        special_group.addWidget(group3_label)
        special_group.addWidget(gov_label)
        special_group.addWidget(self.chk_docnum)
        special_group.addWidget(self.chk_doclevel)
        special_group.addWidget(self.chk_docref)
        special_group.addWidget(med_label)
        special_group.addWidget(self.chk_medicare)
        special_group.addWidget(self.chk_medical_record)
        special_group.addStretch()
        
        types_layout.addLayout(personal_group)
        types_layout.addLayout(business_group)
        types_layout.addLayout(special_group)
        
        scroll_layout.addLayout(types_layout)
        
        # 创建滚动区域
        scroll_area = QScrollArea()
        scroll_area.setWidget(scroll_content)
        scroll_area.setWidgetResizable(True)
        # 让内容自然展开，不限制高度
        # 样式由 theme_manager 统一设置
        
        layout.addWidget(scroll_area)
        
        # 全选/取消按钮
        btn_layout = QHBoxLayout()
        btn_select_all = QPushButton("全选")
        btn_select_all.setObjectName("secondary")
        btn_select_all.clicked.connect(self.select_all_types)
        btn_deselect_all = QPushButton("取消全选")
        btn_deselect_all.setObjectName("secondary")
        btn_deselect_all.clicked.connect(self.deselect_all_types)
        btn_layout.addStretch()
        btn_layout.addWidget(btn_select_all)
        btn_layout.addWidget(btn_deselect_all)
        layout.addLayout(btn_layout)
        
        return frame
    
    def create_content_section(self):
        frame = QFrame()
        frame.setObjectName("card")
        layout = QVBoxLayout(frame)
        layout.setContentsMargins(20, 16, 20, 16)
        
        title = QLabel("文件内容")
        title.setStyleSheet("font-weight: 600; font-size: 14px;")
        
        self.text_edit = QTextEdit()
        self.text_edit.setPlaceholderText("选择文件后，内容将显示在这里...")
        
        # 检测结果标签
        self.result_label = QLabel("")
        self.result_label.setObjectName("status")
        self.result_label.setWordWrap(True)
        
        layout.addWidget(title)
        layout.addWidget(self.text_edit)
        layout.addWidget(self.result_label)
        
        return frame
    
    def create_actions_section(self):
        frame = QFrame()
        layout = QHBoxLayout(frame)
        layout.setContentsMargins(0, 0, 0, 0)

        self.btn_detect = QPushButton("检测敏感信息")
        self.btn_detect.setEnabled(False)
        self.btn_detect.clicked.connect(self.detect_sensitive)

        self.btn_sanitize = QPushButton("开始脱敏")
        self.btn_sanitize.setEnabled(False)
        self.btn_sanitize.clicked.connect(self.sanitize_file)

        self.btn_undo = QPushButton("撤销")
        self.btn_undo.setObjectName("danger")
        self.btn_undo.setEnabled(False)
        self.btn_undo.clicked.connect(self.undo_process)

        self.btn_save = QPushButton("保存结果")
        self.btn_save.setObjectName("secondary")
        self.btn_save.setEnabled(False)
        self.btn_save.clicked.connect(self.save_result)

        layout.addStretch()
        layout.addWidget(self.btn_undo)
        layout.addWidget(self.btn_save)
        layout.addWidget(self.btn_detect)
        layout.addWidget(self.btn_sanitize)

        return frame

    def _apply_scene(self, scene: str):
        """根据场景预设联动勾选——与 batch_tab / settings_tab 完全一致"""
        if not self._scene_ready:
            return
        self._scene_ready = False

        all_checks = [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                       self.chk_ip, self.chk_passport, self.chk_mac, self.chk_imei,
                       self.chk_plate, self.chk_social,
                       self.chk_credit, self.chk_license, self.chk_contract,
                       self.chk_amount, self.chk_phone2,
                       self.chk_docnum, self.chk_doclevel, self.chk_docref,
                       self.chk_medicare, self.chk_medical_record,
                       self.chk_account_permit, self.chk_purchase_order,
                       self.chk_fax, self.chk_employee_id,
                       self.chk_project_code, self.chk_postal]
        for c in all_checks:
            c.setChecked(False)

        general_items = [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                         self.chk_ip, self.chk_passport, self.chk_mac, self.chk_imei,
                         self.chk_plate,
                         self.chk_credit, self.chk_contract, self.chk_amount,
                         self.chk_license, self.chk_phone2,
                         self.chk_account_permit, self.chk_purchase_order,
                         self.chk_fax, self.chk_employee_id,
                         self.chk_project_code, self.chk_postal]

        if scene == 'general':
            for c in general_items:
                c.setChecked(True)
        elif scene == 'gov':
            for c in [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                      self.chk_ip, self.chk_passport, self.chk_plate,
                      self.chk_credit, self.chk_contract, self.chk_amount, self.chk_license,
                      self.chk_docnum, self.chk_doclevel, self.chk_docref,
                      self.chk_account_permit, self.chk_purchase_order,
                      self.chk_fax, self.chk_employee_id,
                      self.chk_project_code, self.chk_postal]:
                c.setChecked(True)
        elif scene == 'finance':
            for c in [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                      self.chk_ip, self.chk_passport, self.chk_mac, self.chk_imei,
                      self.chk_plate, self.chk_social,
                      self.chk_credit, self.chk_contract, self.chk_amount, self.chk_license,
                      self.chk_phone2,
                      self.chk_account_permit, self.chk_purchase_order,
                      self.chk_fax, self.chk_employee_id,
                      self.chk_project_code, self.chk_postal]:
                c.setChecked(True)
        elif scene == 'medical':
            for c in [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                      self.chk_ip, self.chk_passport,
                      self.chk_medicare, self.chk_medical_record,
                      self.chk_social, self.chk_credit, self.chk_contract, self.chk_amount,
                      self.chk_license, self.chk_phone2,
                      self.chk_account_permit, self.chk_purchase_order,
                      self.chk_fax, self.chk_employee_id,
                      self.chk_project_code, self.chk_postal]:
                c.setChecked(True)
        elif scene == 'edu':
            for c in [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                      self.chk_ip, self.chk_passport,
                      self.chk_mac, self.chk_imei, self.chk_plate,
                      self.chk_credit, self.chk_contract, self.chk_amount,
                      self.chk_license, self.chk_phone2,
                      self.chk_account_permit, self.chk_purchase_order,
                      self.chk_fax, self.chk_employee_id,
                      self.chk_project_code, self.chk_postal]:
                c.setChecked(True)
        elif scene == 'custom':
            # 自定义场景：全部勾选
            for c in [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                      self.chk_ip, self.chk_passport, self.chk_mac, self.chk_imei,
                      self.chk_plate, self.chk_social,
                      self.chk_credit, self.chk_license, self.chk_contract,
                      self.chk_amount, self.chk_phone2,
                      self.chk_docnum, self.chk_doclevel, self.chk_docref,
                      self.chk_medicare, self.chk_medical_record,
                      self.chk_account_permit, self.chk_purchase_order,
                      self.chk_fax, self.chk_employee_id,
                      self.chk_project_code, self.chk_postal]:
                c.setChecked(True)

        self._scene_ready = True

    def refresh_sanitize_types(self, types):
        """从设置标签页同步脱敏类型勾选状态"""
        type_map = {
            '手机号': 'chk_phone', '邮箱': 'chk_email', '身份证': 'chk_idcard',
            '银行卡': 'chk_bank', 'IP地址': 'chk_ip', '护照号': 'chk_passport',
            'Mac地址': 'chk_mac', 'IMEI': 'chk_imei', '车牌号': 'chk_plate',
            '社保卡号': 'chk_social', '社会信用代码': 'chk_credit',
            '合同编号': 'chk_contract', '投标/成交价': 'chk_amount',
            '营业执照号': 'chk_license', '固定电话': 'chk_phone2',
            '公文份号': 'chk_docnum', '公文密级': 'chk_doclevel',
            '公文文号': 'chk_docref', '医保号': 'chk_medicare',
            '病历号': 'chk_medical_record', '开户许可证号': 'chk_account_permit',
            '采购/订单编号': 'chk_purchase_order', '传真号': 'chk_fax',
            '工号/学号': 'chk_employee_id', '项目代号': 'chk_project_code',
            '邮编': 'chk_postal',
        }
        for attr in type_map.values():
            chk = getattr(self, attr, None)
            if chk:
                chk.setChecked(False)
        for t in types:
            attr = type_map.get(t)
            if attr:
                chk = getattr(self, attr, None)
                if chk:
                    chk.setChecked(True)

    def select_all_types(self):
        for chk in [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                    self.chk_ip, self.chk_passport, self.chk_mac, self.chk_imei,
                    self.chk_plate, self.chk_social,
                    self.chk_credit, self.chk_license, self.chk_contract,
                    self.chk_amount, self.chk_phone2,
                    self.chk_docnum, self.chk_doclevel, self.chk_docref,
                    self.chk_medicare, self.chk_medical_record,
                    self.chk_account_permit, self.chk_purchase_order,
                    self.chk_fax, self.chk_employee_id,
                    self.chk_project_code, self.chk_postal]:
            chk.setChecked(True)

    def deselect_all_types(self):
        for chk in [self.chk_phone, self.chk_email, self.chk_idcard, self.chk_bank,
                    self.chk_ip, self.chk_passport, self.chk_mac, self.chk_imei,
                    self.chk_plate, self.chk_social,
                    self.chk_credit, self.chk_license, self.chk_contract,
                    self.chk_amount, self.chk_phone2,
                    self.chk_docnum, self.chk_doclevel, self.chk_docref,
                    self.chk_medicare, self.chk_medical_record,
                    self.chk_account_permit, self.chk_purchase_order,
                    self.chk_fax, self.chk_employee_id,
                    self.chk_project_code, self.chk_postal]:
            chk.setChecked(False)

    def undo_process(self):
        """撤销脱敏，恢复原始内容"""
        if not self.current_file or not self.original_content:
            return

        reply = QMessageBox.question(
            self, "确认撤销", "确定要撤销脱敏，恢复原始内容吗？",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )

        if reply == QMessageBox.StandardButton.Yes:
            self.text_edit.setPlainText(self.original_content)
            self.processed_content = None
            self.btn_undo.setEnabled(False)
            self.btn_save.setEnabled(False)
            self.result_label.setText("↩️ 已撤销，恢复原始内容")

    def set_file(self, file_path):
        """设置文件（供拖拽使用）"""
        self.current_file = file_path
        self.file_label.setText(Path(file_path).name)
        self.load_file_content(file_path)
        self.btn_detect.setEnabled(True)
        self.btn_sanitize.setEnabled(True)
        self.btn_undo.setEnabled(False)

    def browse_file(self):
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "选择文件",
            "",
            "支持的文件 (*.txt *.md *.json *.xml *.csv *.html *.docx);;所有文件 (*.*)"
        )
        if file_path:
            self.current_file = file_path
            self.file_label.setText(Path(file_path).name)
            self.load_file_content(file_path)
            self.btn_detect.setEnabled(True)
            self.btn_sanitize.setEnabled(True)
    
    def _get_preview_text(self, path):
        """获取文件的纯文本预览（用于 text_edit 显示）
        
        对于 docx/xlsx/pptx，提取纯文本用于显示；
        对于 txt/md/json 等，直接读取文本。
        格式完全保持的原地脱敏在保存时通过原生对象实现。
        
        返回: (text, native_obj, actual_ext)
        """
        ext = Path(path).suffix.lower()
        actual_ext = ext  # 记录实际处理的格式
        
        # .doc → 先转 .docx 再处理
        if ext == '.doc':
            try:
                from doc2docx import convert as doc2docx_convert
                import tempfile, os
                with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
                    tmp_docx = tmp.name
                try:
                    doc2docx_convert(str(path), tmp_docx)
                except Exception:
                    if os.path.exists(tmp_docx):
                        os.remove(tmp_docx)
                    return self._fallback_read(path), None, ext
                path = tmp_docx
                ext = '.docx'
                actual_ext = '.docx'
            except ImportError:
                return self._fallback_read(path), None, ext
        
        if ext == '.docx':
            try:
                from docx import Document
                doc = Document(path)
                paragraphs = [p.text for p in doc.paragraphs]
                return '\n'.join(paragraphs), doc, actual_ext
            except ImportError:
                return self._fallback_read(path), None, actual_ext
            except Exception as e:
                raise Exception(f"无法读取 docx 文件: {e}")
        
        elif ext in ('.xlsx', '.xls'):
            try:
                from openpyxl import load_workbook
                # data_only=False + read_only=False：显示公式文字（如 =SUM(B2:B5)）
                # 注意：只读模式 data_only 参数无效，必须关闭只读才能看到公式
                wb = load_workbook(path, data_only=False)
                lines = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        cells = []
                        for c in row:
                            if c is None:
                                cells.append('')
                            elif isinstance(c, float) and (c != c):  # NaN check
                                cells.append('')
                            else:
                                cells.append(str(c))
                        row_text = '\t'.join(cells)
                        if row_text.strip():
                            lines.append(row_text)
                # 不要 close()！_native_doc 后续由 _sanitize_native_xlsx/save_result 使用
                return '\n'.join(lines), wb, actual_ext
                return '\n'.join(lines), wb, actual_ext
            except ImportError:
                return self._fallback_read(path), None, actual_ext
            except Exception as e:
                raise Exception(f"无法读取 xlsx 文件: {e}")
        
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(path)
                lines = []
                for i, slide in enumerate(prs.slides, 1):
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    if slide_texts:
                        lines.append(f'[幻灯片 {i}]\n' + '\n'.join(slide_texts))
                return '\n'.join(lines), prs, actual_ext
            except ImportError:
                return self._fallback_read(path), None, actual_ext
            except Exception as e:
                raise Exception(f"无法读取 pptx 文件: {e}")
        
        elif ext == '.pdf':
            # PDF 文件：使用 pypdf 提取文本
            try:
                from pypdf import PdfReader
                reader = PdfReader(path)
                lines = []
                for i, page in enumerate(reader.pages, 1):
                    text = page.extract_text() or ''
                    if text.strip():
                        lines.append(f'[页 {i}]\n{text}')
                return '\n'.join(lines), None, actual_ext
            except ImportError:
                return self._fallback_read(path), None, actual_ext
            except Exception as e:
                raise Exception(f"无法读取 pdf 文件: {e}")
        
        else:
            # txt/md/json/xml/csv/html 等纯文本文档，原生对象即内容本身
            return self._fallback_read(path), None, actual_ext

    def _fallback_read(self, path):
        """回退的文本读取（用于不支持原生处理的格式）"""
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def load_file_content(self, path):
        try:
            content, native_obj, actual_ext = self._get_preview_text(path)
            self.original_content = content
            self.processed_content = None
            self._native_doc = native_obj  # 可能是 Document / Workbook / Presentation 或 None
            self._native_docx_ext = actual_ext  # .doc 转换后为 .docx
            self.text_edit.setPlainText(content)
            self.btn_undo.setEnabled(False)
        except Exception as e:
            QMessageBox.warning(self, "错误", f"无法读取文件: {e}")
    
    def get_selected_types(self):
        """获取选中的脱敏类型"""
        types = []
        if self.chk_phone.isChecked(): types.append('手机号')
        if self.chk_email.isChecked(): types.append('邮箱')
        if self.chk_idcard.isChecked(): types.append('身份证')
        if self.chk_bank.isChecked(): types.append('银行卡')
        if self.chk_ip.isChecked(): types.append('IP地址')
        if self.chk_passport.isChecked(): types.append('护照号')
        if self.chk_mac.isChecked(): types.append('Mac地址')
        if self.chk_imei.isChecked(): types.append('IMEI')
        if self.chk_plate.isChecked(): types.append('车牌号')
        if self.chk_social.isChecked(): types.append('社保卡号')
        if self.chk_credit.isChecked(): types.append('社会信用代码')
        if self.chk_license.isChecked(): types.append('营业执照号')
        if self.chk_contract.isChecked(): types.append('合同编号')
        if self.chk_amount.isChecked(): types.append('投标/成交价')
        if self.chk_phone2.isChecked(): types.append('固定电话')
        if self.chk_docnum.isChecked(): types.append('公文份号')
        if self.chk_doclevel.isChecked(): types.append('公文密级')
        if self.chk_docref.isChecked(): types.append('公文文号')
        if self.chk_medicare.isChecked(): types.append('医保号')
        if self.chk_medical_record.isChecked(): types.append('病历号')
        if self.chk_account_permit.isChecked(): types.append('开户许可证号')
        if self.chk_purchase_order.isChecked(): types.append('采购/订单编号')
        if self.chk_fax.isChecked(): types.append('传真号')
        if self.chk_employee_id.isChecked(): types.append('工号/学号')
        if self.chk_project_code.isChecked(): types.append('项目代号')
        if self.chk_postal.isChecked(): types.append('邮编')
        return types
    
    def get_custom_patterns(self):
        """从设置加载自定义规则"""
        settings_file = Path(__file__).parent / 'settings.json'
        if settings_file.exists():
            try:
                with open(settings_file, 'r', encoding='utf-8') as f:
                    settings = json.load(f)
                return settings.get('custom_patterns', {'keywords': '', 'regex': ''})
            except:
                pass
        return {'keywords': '', 'regex': ''}

    def detect_sensitive(self):
        """检测敏感信息"""
        if not self.current_file:
            return
        
        from safe_shrink_gui import detect_sensitive
        
        content = self.text_edit.toPlainText()
        types = self.get_selected_types()
        custom = self.get_custom_patterns()
        
        self.detected_items = detect_sensitive(content, types, custom_patterns=custom)
        
        if self.detected_items:
            # 统计各类型数量
            type_counts = {}
            for item in self.detected_items:
                t = item['type']
                type_counts[t] = type_counts.get(t, 0) + 1
            
            msg = f"检测到 {len(self.detected_items)} 处敏感信息:\n"
            for t, count in type_counts.items():
                msg += f"- {t}: {count} 处\n"
            
            self.result_label.setText(msg.strip())
            QMessageBox.information(self, "检测结果", msg)
        else:
            self.result_label.setText("未检测到敏感信息")
            QMessageBox.information(self, "检测结果", "未检测到敏感信息")

    def _sanitize_text(self, text, types, custom, mode):
        """对纯文本执行脱敏，返回脱敏后的文本"""
        import sys, traceback
        try:
            from safe_shrink_gui import sanitize_content
            return sanitize_content(text, types, custom_patterns=custom, mode=mode)
        except ImportError as e:
            # 详细调试信息
            print(f"[DEBUG] ImportError: {e}")
            print(f"[DEBUG] sys.frozen: {getattr(sys, 'frozen', False)}")
            print(f"[DEBUG] sys._MEIPASS: {getattr(sys, '_MEIPASS', 'N/A')}")
            print(f"[DEBUG] sys.path: {sys.path[:3]}...")
            if 'safe_shrink_gui' in sys.modules:
                mod = sys.modules['safe_shrink_gui']
                print(f"[DEBUG] safe_shrink_gui in sys.modules: {mod}")
                print(f"[DEBUG] dir(mod): {[x for x in dir(mod) if not x.startswith('_')][:10]}")
            else:
                print(f"[DEBUG] safe_shrink_gui NOT in sys.modules")
            traceback.print_exc()
            raise
    
    def _sanitize_native_docx(self, doc, types, custom, mode):
        """原地脱敏 docx Document 对象（保持所有格式）"""
        import sys, traceback
        print(f"[DEBUG] _sanitize_native_docx called, types={types}")
        try:
            from safe_shrink_gui import sanitize_content
        except ImportError as e:
            print(f"[DEBUG] _sanitize_native_docx ImportError: {e}")
            if 'safe_shrink_gui' in sys.modules:
                mod = sys.modules['safe_shrink_gui']
                print(f"[DEBUG] safe_shrink_gui in sys.modules: {mod}")
                attrs = [x for x in dir(mod) if not x.startswith('_')]
                print(f"[DEBUG] attrs: {attrs[:15]}")
            else:
                print(f"[DEBUG] safe_shrink_gui NOT in sys.modules")
            traceback.print_exc()
            raise
        import re
        
        masked = {}
        def replacer(match):
            placeholder = f'__MASK_{len(masked)}__'
            masked[placeholder] = match.group(0)
            return placeholder
        
        def sanitize_run_text(run_text):
            if not run_text.strip():
                return run_text
            # 先用占位符保护已有的遮蔽内容
            protected = re.sub(r'__MASK_\d+__', replacer, run_text)
            # 执行脱敏
            sanitized = sanitize_content(protected, types, custom_patterns=custom, mode=mode)
            # 还原占位符
            for ph, original in masked.items():
                sanitized = sanitized.replace(ph, original)
            return sanitized
        
        # 收集所有段落（含表格内）
        all_paragraphs = list(doc.paragraphs)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    all_paragraphs.extend(cell.paragraphs)
        
        # 保存原始段落文本（用于跨 run 检测）
        original_texts = {}
        for para in all_paragraphs:
            original_texts[id(para)] = para.text
        
        # Pass 1: 逐 run 脱敏（保留格式）
        for para in all_paragraphs:
            for run in para.runs:
                new_text = sanitize_run_text(run.text)
                if new_text != run.text:
                    run.text = new_text
        
        # Pass 2: 跨 run 修复 — 检查段落合并文本是否仍有残留敏感数据
        for para in all_paragraphs:
            original = original_texts[id(para)]
            if not original.strip():
                continue
            # 对原始段落文本做完整脱敏
            fully_sanitized = sanitize_content(original, types, custom_patterns=custom, mode=mode)
            current = para.text
            if fully_sanitized != current and para.runs:
                # 存在跨 run 遗漏，合并到第一个 run
                para.runs[0].text = fully_sanitized
                for run in para.runs[1:]:
                    run.text = ''
    
    def _sanitize_native_xlsx(self, wb, types, custom, mode):
        """原地脱敏 xlsx（保持所有格式）
        注意：传入的 wb 可能是 read_only=True 模式（来自预览），不可修改。
        此函数用 read_only=False 重新加载一份可写副本。
        """
        import sys, traceback
        try:
            from openpyxl import load_workbook as _load
            from safe_shrink_gui import sanitize_content
        except ImportError as e:
            print(f"[DEBUG] _sanitize_native_xlsx ImportError: {e}")
            traceback.print_exc()
            raise
        
        # 读取当前文件路径，用可写模式重新加载
        file_path = getattr(self, '_current_file_path', None) or getattr(self, 'current_file', '')
        if not file_path:
            # fallback: 尝试关闭只读 wb，用同一路径重新打开
            try:
                wb.close()
            except Exception:
                pass
            raise Exception("无法确定 xlsx 文件路径，请重新打开文件")
        
        # 保存预览用的只读 workbook（如果可迭代）
        try:
            wb.close()
        except Exception:
            pass
        
        # 用可写模式加载
        writable_wb = _load(str(file_path), data_only=False)  # False: 保留公式，Excel打开后重算
        self._native_doc = writable_wb  # 更新引用为可写版本
        
        for sheet_name in writable_wb.sheetnames:
            ws = writable_wb[sheet_name]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        cell_text = str(cell.value)
                        sanitized = sanitize_content(cell_text, types, custom_patterns=custom, mode=mode)
                        if sanitized != cell_text:
                            cell.value = sanitized
    
    def _sanitize_native_pptx(self, prs, types, custom, mode):
        """原地脱敏 python-pptx Presentation 对象（保持所有格式）"""
        import sys, traceback
        try:
            from safe_shrink_gui import sanitize_content
        except ImportError as e:
            print(f"[DEBUG] _sanitize_native_pptx ImportError: {e}")
            traceback.print_exc()
            raise
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            run.text = sanitize_content(
                                run.text, types, custom_patterns=custom, mode=mode
                            )
    
    def _update_preview_from_native(self):
        """在原地脱敏后，用原生文档更新预览文本"""
        ext = self._native_docx_ext
        
        if ext == '.docx' and self._native_doc is not None:
            paras = [p.text for p in self._native_doc.paragraphs]
            self.text_edit.setPlainText('\n'.join(paras))
        
        elif ext in ('.xlsx', '.xls') and self._native_doc is not None:
            lines = []
            for sheet_name in self._native_doc.sheetnames:
                ws = self._native_doc[sheet_name]
                for row in ws.iter_rows(values_only=True):
                    cells = []
                    for c in row:
                        if c is None:
                            cells.append('')
                        elif isinstance(c, float) and (c != c):  # NaN check
                            cells.append('')
                        else:
                            cells.append(str(c))
                    row_text = '\t'.join(cells)
                    if row_text.strip():
                        lines.append(row_text)
            # preview 函数里 _native_doc 来自 _get_preview_text() 的 read_only=True 模式
            # preview 不关闭它（由 _update_preview_from_native 统一管理生命周期）
            self.text_edit.setPlainText('\n'.join(lines))
        
        elif ext == '.pptx' and self._native_doc is not None:
            lines = []
            for i, slide in enumerate(self._native_doc.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    lines.append(f'[幻灯片 {i}]\n' + '\n'.join(slide_texts))
            self.text_edit.setPlainText('\n'.join(lines))
    
    def sanitize_file(self):
        """执行脱敏（原地处理，保持格式）"""
        if not self.current_file:
            return

        types = self.get_selected_types()
        print(f"[DEBUG] sanitize_file called, types={types}")  # 调试输出
        custom = self.get_custom_patterns()
        mode = 'mask' if self.radio_mask.isChecked() else 'tags'

        try:
            # 如果有原生文档对象，优先原地处理
            if self._native_doc is not None:
                ext = self._native_docx_ext
                
                if ext == '.docx':
                    self._sanitize_native_docx(self._native_doc, types, custom, mode)
                    self._update_preview_from_native()
                
                elif ext in ('.xlsx', '.xls'):
                    self._sanitize_native_xlsx(self._native_doc, types, custom, mode)
                    self._update_preview_from_native()
                
                elif ext == '.pptx':
                    self._sanitize_native_pptx(self._native_doc, types, custom, mode)
                    self._update_preview_from_native()
                
                self.processed_content = self.text_edit.toPlainText()
            
            else:
                # 纯文本文件：直接用 safe_shrink_gui 脱敏
                content = self.text_edit.toPlainText()
                result = self._sanitize_text(content, types, custom, mode)
                self.processed_content = result
                self.text_edit.setPlainText(result)

            self.btn_save.setEnabled(True)
            self.btn_undo.setEnabled(True)
            self.result_label.setText("✅ 脱敏完成！格式已保留。")
            QMessageBox.information(self, "完成", "脱敏完成！格式已保留。")
        except Exception as e:
            QMessageBox.critical(self, "错误", f"处理失败: {e}")

    def save_result(self):
        """保存脱敏结果（使用原生对象保存，完全保持格式）"""
        if not self.current_file:
            return

        original_path = Path(self.current_file)
        original_ext = self._native_docx_ext or original_path.suffix.lower()
        original_file_ext = original_path.suffix.lower()  # 用户原始文件扩展名
        
        # 根据原文件类型设置默认保存路径和过滤器
        # 注意：.doc 文件内部转换为 .docx 处理，保存时也只能保存为 .docx
        if original_ext == '.docx':
            default_name = str(original_path.with_stem(original_path.stem + '_脱敏'))
            file_filter = "Word文档 (*.docx)"
            if original_file_ext == '.doc':
                # 原 .doc 文件，提示将保存为 .docx
                file_filter = "Word文档 (*.docx) - 原文件为.doc，自动转为.docx格式"
        elif original_ext in ('.xlsx', '.xls'):
            default_name = str(original_path.with_stem(original_path.stem + '_脱敏'))
            file_filter = "Excel表格 (*.xlsx)"
        elif original_ext == '.pptx':
            default_name = str(original_path.with_stem(original_path.stem + '_脱敏'))
            file_filter = "PowerPoint演示文稿 (*.pptx)"
        elif original_ext == '.pdf':
            default_name = str(original_path.with_stem(original_path.stem + '_脱敏').with_suffix('.txt'))
            file_filter = "文本文件 (*.txt);;Markdown文件 (*.md)"
        elif original_ext in ('.ssd', '.md'):
            default_name = str(original_path.with_stem(original_path.stem + '_脱敏'))
            file_filter = "Markdown文件 (*.md);;SSD文档 (*.ssd);;文本文件 (*.txt)"
        else:
            default_name = str(original_path.with_suffix('.sanitized.txt'))
            file_filter = "文本文件 (*.txt)"

        file_path, _ = QFileDialog.getSaveFileName(
            self,
            "保存文件",
            default_name,
            file_filter
        )
        
        if file_path:
            try:
                save_ext = Path(file_path).suffix.lower()
                
                # .doc 格式不支持，强制转为 .docx
                if save_ext == '.doc' and original_ext == '.docx':
                    file_path = str(Path(file_path).with_suffix('.docx'))
                    save_ext = '.docx'
                    QMessageBox.information(self, "格式转换", 
                        ".doc 格式不支持，已自动保存为 .docx 格式")
                
                # PDF 脱敏后无法保存回 PDF，提示用户
                if original_ext == '.pdf' and save_ext == '.pdf':
                    file_path = str(Path(file_path).with_suffix('.txt'))
                    save_ext = '.txt'
                    QMessageBox.information(self, "格式转换", 
                        "PDF 脱敏后无法保持原格式，已自动保存为文本文件")
                
                # 优先使用原生文档对象保存（保持格式）
                if self._native_doc is not None and save_ext == original_ext:
                    if save_ext == '.docx':
                        self._native_doc.save(file_path)
                    elif save_ext in ('.xlsx', '.xls'):
                        self._native_doc.save(file_path)
                    elif save_ext == '.pptx':
                        self._native_doc.save(file_path)
                
                # 格式不同时回退到纯文本保存
                else:
                    content = self.text_edit.toPlainText()
                    if save_ext == '.docx':
                        self._save_as_docx(file_path, content)
                    elif save_ext in ('.xlsx', '.xls'):
                        self._save_as_xlsx(file_path, content)
                    elif save_ext == '.pptx':
                        self._save_as_pptx(file_path, content)
                    else:
                        with open(file_path, 'w', encoding='utf-8') as f:
                            f.write(content)
                
                QMessageBox.information(self, "保存成功", f"已保存到: {file_path}")
            except Exception as e:
                QMessageBox.critical(self, "错误", f"保存失败: {e}")
    
    def _save_as_docx(self, filepath, content):
        """保存为Word文档"""
        try:
            from docx import Document
        except ImportError:
            # 回退到 safe_shrink 的写入函数
            from safe_shrink import write_file
            write_file(filepath, content, 'docx')
            return
        
        doc = Document()
        for para in content.split('\n'):
            if para.strip():
                doc.add_paragraph(para)
        doc.save(filepath)
    
    def _save_as_xlsx(self, filepath, content):
        """保存为Excel表格"""
        try:
            from openpyxl import Workbook
        except ImportError:
            from safe_shrink import write_file
            write_file(filepath, content, 'xlsx')
            return
        
        wb = Workbook()
        ws = wb.active
        for i, line in enumerate(content.split('\n'), 1):
            if '\t' in line:
                for j, cell in enumerate(line.split('\t')):
                    ws.cell(row=i, column=j+1, value=cell)
            elif line.strip():
                ws.cell(row=i, column=1, value=line)
        wb.save(filepath)
    
    def _save_as_pptx(self, filepath, content):
        """保存为PowerPoint演示文稿"""
        from safe_shrink import write_file
        write_file(filepath, content, 'pptx')

    def cleanup(self):
        pass
    
    def update_language(self, lang):
        """更新语言"""
        from translations import get_translation
        _ = lambda t: get_translation(t, lang)
        
        # 更新场景预设
        if hasattr(self, 'scene_general'):
            self.scene_general.setText(_('通用文档'))
        if hasattr(self, 'scene_gov'):
            self.scene_gov.setText(_('党政公文'))
        if hasattr(self, 'scene_finance'):
            self.scene_finance.setText(_('金融合同'))
        if hasattr(self, 'scene_medical'):
            self.scene_medical.setText(_('医疗档案'))
        if hasattr(self, 'scene_edu'):
            self.scene_edu.setText(_('教育材料'))
        if hasattr(self, 'scene_custom'):
            self.scene_custom.setText(_('自定义'))
        
        # 更新脱敏模式
        if hasattr(self, 'radio_mask'):
            self.radio_mask.setText(_('简单遮蔽 (***)'))
        if hasattr(self, 'radio_tags'):
            self.radio_tags.setText(_('语义标签 (<PERSON_NAME.1>)'))
        
        # 更新按钮
        if hasattr(self, 'btn_detect'):
            self.btn_detect.setText(_('检测敏感信息'))
        if hasattr(self, 'btn_sanitize'):
            self.btn_sanitize.setText(_('开始脱敏'))
        if hasattr(self, 'btn_undo'):
            self.btn_undo.setText(_('撤销'))
        if hasattr(self, 'btn_save'):
            self.btn_save.setText(_('保存结果'))
        
        # 更新全选按钮
        # These are dynamically created, need to find them by text
        pass
