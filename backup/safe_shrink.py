# -*- coding: utf-8 -*-

import sys as _sys

try:

    _sys.stdout.reconfigure(encoding='utf-8')

    _sys.stderr.reconfigure(encoding='utf-8')

except Exception:

    pass



"""

SafeShrink 增强版 v2.0

=====================

支持格式:

  - 纯文本: .txt, .md, .json, .csv, .xml, .html, .log

  - Word: .docx (需要 python-docx)

  - Excel: .xlsx (需要 openpyxl, 已有), .xls (需要 xlrd)

  - PPT: .pptx (需要 python-pptx, 已有)

  - PDF: .pdf (需要 pdfplumber 或 PyMuPDF)



用法:

  python safe_shrink.py slim input.docx -o output.txt

  python safe_shrink.py sanitize input.pdf -o output.pdf

  python safe_shrink.py slim input.xlsx --sheet 0

  python safe_shrink.py sanitize input.docx -o sanitized.docx --format docx

"""



import argparse

import json

import sys

from datetime import datetime

import os

import re

import zipfile

import csv

from pathlib import Path



# ========== 依赖检测 ==========

DEPS = {}

try:

    import docx

    DEPS['docx'] = True

except ImportError:

    DEPS['docx'] = False



try:

    import openpyxl

    DEPS['openpyxl'] = True

except ImportError:

    DEPS['openpyxl'] = False



try:

    import pptx

    DEPS['pptx'] = True

except ImportError:

    DEPS['pptx'] = False



try:

    import pdfplumber

    DEPS['pdfplumber'] = True

except ImportError:

    DEPS['pdfplumber'] = False



try:

    from pypdf import PdfReader, PdfWriter

    DEPS['pypdf'] = True

except ImportError:

    DEPS['pypdf'] = False



try:

    import xlrd

    DEPS['xlrd'] = True

except ImportError:

    DEPS['xlrd'] = False



try:

    from PIL import Image

    DEPS['pillow'] = True

except ImportError:

    DEPS['pillow'] = False





def check_dep(name, pkg_name=None):

    """检查依赖，打印警告"""

    if DEPS.get(name):

        return True

    pkg = pkg_name or name

    print(f"[警告] 缺少 {pkg}，相关功能不可用。安装: pip install {pkg}", file=sys.stderr)

    return False





# ========== 文件读取模块 ==========



def read_txt(filepath):

    """读取纯文本文件"""

    for enc in ['utf-8', 'gbk', 'gb2312', 'utf-16']:

        try:

            with open(filepath, 'r', encoding=enc) as f:

                return f.read()

        except:

            continue

    raise ValueError(f"无法读取文件: {filepath}")





def read_json(filepath):

    """读取JSON文件，返回原始文本"""

    with open(filepath, 'r', encoding='utf-8') as f:

        obj = json.load(f)

    return json.dumps(obj, ensure_ascii=False, indent=2)





def read_csv(filepath):

    """读取CSV文件"""

    lines = []

    with open(filepath, 'r', encoding='utf-8-sig', newline='') as f:

        reader = csv.reader(f)

        for row in reader:

            lines.append(','.join(row))

    return '\n'.join(lines)





def read_docx(filepath):

    """读取Word .docx文件"""

    if not check_dep('docx', 'python-docx'):

        raise ImportError("python-docx 未安装")

    import docx

    doc = docx.Document(filepath)

    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # 读表格

    tables_text = []

    for table in doc.tables:

        for row in table.rows:

            cells = [cell.text.strip() for cell in row.cells]

            if any(cells):

                tables_text.append(' | '.join(cells))

    result = '\n'.join(paragraphs)

    if tables_text:

        result += '\n\n[表格]\n' + '\n'.join(tables_text)

    return result





def read_xlsx(filepath, sheet_index=0):

    """读取Excel .xlsx文件"""

    if not check_dep('openpyxl'):

        raise ImportError("openpyxl 未安装")

    import openpyxl

    wb = openpyxl.load_workbook(filepath, data_only=True)

    sheets = wb.sheetnames

    result_parts = []

    for idx in (sheet_index if isinstance(sheet_index, list) else [sheet_index]):

        if idx < len(sheets):

            ws = wb[sheets[idx]]

            result_parts.append(f"[Sheet: {sheets[idx]}]")

            for row in ws.iter_rows(values_only=True):

                cells = [str(c) if c is not None else '' for c in row]

                if any(cells):

                    result_parts.append('\t'.join(cells))

    return '\n'.join(result_parts)





def read_xls(filepath, sheet_index=0):

    """读取Excel .xls文件（老格式）"""

    if not check_dep('xlrd', 'xlrd'):

        raise ImportError("xlrd 未安装 (需要支持.xls)")

        # 退而求其次: 尝试openpyxl

        if check_dep('openpyxl'):

            try:

                import openpyxl

                wb = openpyxl.load_workbook(filepath)

                ws = wb.active

                lines = []

                for row in ws.iter_rows(values_only=True):

                    cells = [str(c) if c else '' for c in row]

                    if any(cells):

                        lines.append('\t'.join(cells))

                return '\n'.join(lines)

            except:

                pass

        raise ImportError("xlrd 未安装，无法读取.xls文件")

    import xlrd

    wb = xlrd.open_workbook(filepath)

    sheets = wb.sheet_names()

    parts = []

    idx = sheet_index if isinstance(sheet_index, list) else sheet_index

    if idx < len(sheets):

        ws = wb.sheet_by_index(idx)

        parts.append(f"[Sheet: {sheets[idx]}]")

        for r in range(ws.nrows):

            row_data = [str(ws.cell_value(r, c)) for c in range(ws.ncols)]

            if any(c.strip() for c in row_data):

                parts.append('\t'.join(row_data))

    return '\n'.join(parts)





def read_pptx(filepath):

    """读取PowerPoint .pptx文件"""

    if not check_dep('pptx', 'python-pptx'):

        raise ImportError("python-pptx 未安装")

    import pptx

    prs = pptx.Presentation(filepath)

    slides_text = []

    for i, slide in enumerate(prs.slides, 1):

        slide_texts = []

        for shape in slide.shapes:

            if hasattr(shape, 'text') and shape.text.strip():

                slide_texts.append(shape.text.strip())

        if slide_texts:

            slides_text.append(f"[幻灯片 {i}]\n" + '\n'.join(slide_texts))

    return '\n\n'.join(slides_text)





def read_pdf_pdfplumber(filepath):

    """用pdfplumber读取PDF"""

    if not check_dep('pdfplumber', 'pdfplumber'):

        raise ImportError("pdfplumber 未安装")

    import pdfplumber

    pages_text = []

    with pdfplumber.open(filepath) as pdf:

        for i, page in enumerate(pdf.pages, 1):

            text = page.extract_text()

            if text and text.strip():

                pages_text.append(f"[页 {i}]\n{text}")

    return '\n\n'.join(pages_text)





def read_pdf_pypdf(filepath):

    """用pypdf读取PDF"""

    if not check_dep('pypdf', 'pypdf'):

        raise ImportError("pypdf 未安装")

    from pypdf import PdfReader

    reader = PdfReader(filepath)

    pages_text = []

    for i, page in enumerate(reader.pages, 1):

        text = page.extract_text()

        if text and text.strip():

            pages_text.append(f"[页 {i}]\n{text}")

    return '\n\n'.join(pages_text)





def read_pdf(filepath):

    """读取PDF，尝试多种方式"""

    errors = []

    # 优先 pdfplumber（表格支持好）

    if DEPS.get('pdfplumber'):

        try:

            return read_pdf_pdfplumber(filepath)

        except Exception as e:

            errors.append(f"pdfplumber: {e}")

    # 备选 pypdf

    if DEPS.get('pypdf'):

        try:

            return read_pdf_pypdf(filepath)

        except Exception as e:

            errors.append(f"pypdf: {e}")

    raise ImportError(f"无法读取PDF。尝试了: {', '.join(errors)}")





# ========== 文件写入模块 ==========



def write_txt(filepath, content):

    with open(filepath, 'w', encoding='utf-8') as f:

        f.write(content)





def write_json(filepath, content):

    obj = json.loads(content)

    with open(filepath, 'w', encoding='utf-8') as f:

        json.dump(obj, f, ensure_ascii=False, indent=2)





def write_docx(filepath, content):

    """写入Word .docx"""

    if not check_dep('docx', 'python-docx'):

        raise ImportError("python-docx 未安装")

    import docx

    doc = docx.Document()

    for para in content.split('\n'):

        if para.strip():

            doc.add_paragraph(para)

    doc.save(filepath)





def write_xlsx(filepath, content):

    """写入Excel .xlsx"""

    if not check_dep('openpyxl'):

        raise ImportError("openpyxl 未安装")

    import openpyxl

    wb = openpyxl.Workbook()

    ws = wb.active

    for i, line in enumerate(content.split('\n'), 1):

        if '\t' in line:

            for j, cell in enumerate(line.split('\t')):

                ws.cell(row=i, column=j+1, value=cell)

        elif line.strip():

            ws.cell(row=i, column=1, value=line)

    wb.save(filepath)





def write_pptx(filepath, content):

    """写入PowerPoint .pptx"""

    if not check_dep('pptx', 'python-pptx'):

        raise ImportError("python-pptx 未安装")

    import pptx

    prs = pptx.Presentation()

    for block in content.split('[幻灯片 '):

        if not block.strip():

            continue

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局

        lines = block.split('\n')

        title = lines[0].rstrip(']') if lines else f"内容 {len(prs.slides)}"

        body = '\n'.join(lines[1:]) if len(lines) > 1 else block

        if hasattr(slide.shapes, 'title') and slide.shapes.title:

            slide.shapes.title.text = title

        tf = slide.shapes.placeholders[0].text_frame if slide.shapes.placeholders else None

        if tf:

            tf.text = body

        else:

            from pptx.util import Inches, Pt

            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))

            tf2 = txBox.text_frame

            tf2.text = body

    prs.save(filepath)





def write_pdf_pdfplumber(filepath, content):

    """用pdfplumber写入PDF（实际上是文本到PDF）"""

    try:

        from reportlab.pdfgen import canvas

        from reportlab.lib.pagesizes import A4

        from reportlab.pdfbase import pdfmetrics

        from reportlab.pdfbase.ttfonts import TTFont

        from reportlab.lib.styles import getSampleStyleSheet

        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    except ImportError:

        # 没有reportlab，用备选方案：保存为同名txt

        alt = str(Path(filepath).with_suffix('.txt'))

        with open(alt, 'w', encoding='utf-8') as f:

            f.write(content)

        raise RuntimeError(f"PDF写入需要reportlab，已保存为: {alt}")

    import pdfplumber

    pdf = pdfplumber.open(filepath)

    # 实际pdfplumber不支持创建PDF，这里只是保留接口

    raise RuntimeError("pdfplumber不支持创建PDF，请使用文本输出")





# ========== 文件读取入口 ==========



READERS = {

    '.txt': read_txt,

    '.md': read_txt,

    '.ssd': read_txt,

    '.json': read_json,

    '.csv': read_csv,

    '.xml': read_txt,

    '.html': read_txt,

    '.htm': read_txt,

    '.log': read_txt,

    '.ini': read_txt,

    '.cfg': read_txt,

    '.yaml': read_txt,

    '.yml': read_txt,

    '.docx': read_docx,

    '.xlsx': read_xlsx,

    '.xls': read_xls,

    '.pptx': read_pptx,

    '.pdf': read_pdf,

}





def read_file(filepath, options=None):

    """统一读取接口"""

    opts = options or {}

    ext = Path(filepath).suffix.lower()

    

    if ext == '.xlsx' and 'sheet' in opts:

        return read_xlsx(filepath, sheet_index=opts['sheet'])

    if ext == '.xls' and 'sheet' in opts:

        return read_xls(filepath, sheet_index=opts['sheet'])

    

    if ext not in READERS:

        raise ValueError(f"不支持的格式: {ext}。支持: {', '.join(READERS.keys())}")

    

    return READERS[ext](filepath)





def write_file(filepath, content, fmt=None):

    """统一写入接口"""

    ext = Path(filepath).suffix.lower() if fmt is None else f'.{fmt}'

    

    writers = {

        '.txt': write_txt,

        '.md': write_txt,

        '.ssd': write_txt,

        '.json': write_json,

        '.csv': write_txt,

        '.docx': write_docx,

        '.xlsx': write_xlsx,

        '.pptx': write_pptx,

    }

    

    if ext not in writers:

        # 默认写txt

        ext = '.txt'

        alt_path = str(Path(filepath).with_suffix('.txt'))

        print(f"[注意] 此格式不支持直接写入，已保存为: {alt_path}", file=sys.stderr)

        filepath = alt_path

    

    writers[ext](filepath, content)





# ========== Token 估算 ==========



import re as _re



def estimate_tokens(text):

    """

    估算文本的 Token 消耗（基于 GPT-4o / Claude 等主流大模型）

    

    参数:

        text (str): 输入文本（纯文本或 Markdown）

    

    返回:

        dict: {

            'total': int,       # 总 token 数

            'chinese': int,     # 中文 token

            'english': int,     # 英文 token

            'numbers': int,     # 数字 token

            'punctuation': int, # 标点 token

            'whitespace': int,  # 空白 token

            'tables': int,      # 表格 token

            'code_blocks': int, # 代码块 token

            'links': int,       # 链接 token

            'images': int,      # 图片 token（仅 Markdown 图片标记）

            'lists': int,       # 列表 token

            'headings': int,    # 标题 token

            'emphasis': int,    # 加粗/斜体 token

            'quotes': int,      # 引用块 token

            'hr': int,          # 水平线 token

            'footnotes': int,   # 脚注 token

            'math': int,        # 数学公式 token

            'html': int,        # HTML 标签 token

            'strikethrough': int,  # 删除线 token

            'checkboxes': int,  # Checkbox token

            'toc': int,         # 目录 token

        }

    """

    if not text:

        return _empty_token_result()

    

    stats = _empty_token_result()

    

    # --- 先统计 Markdown 结构元素（这些不参与基础字符统计）---

    

    # 1. 图片: ![alt](url) — 图片本身按 ~170 tokens 估算（中等分辨率）

    for m in _re.finditer(r'!\[([^\]]*)\]\(([^)]+)\)', text):

        stats['images'] += 170

    

    # 2. 超链接: [text](url) — 非图片

    link_matches = _re.findall(r'(?<!!)\[([^\]]+)\]\(([^)]+)\)', text)

    stats['links'] = len(link_matches) * 15

    

    # 3. 代码块: ```...```

    code_blocks = _re.findall(r'```[^\n]*\n(.*?)```', text, _re.DOTALL)

    for block in code_blocks:

        stats['code_blocks'] += 20  # 开销

        stats['code_blocks'] += block.count('\n') * 3  # 每行

    # 行内代码: `code`

    inline_codes = _re.findall(r'`([^`]+)`', text)

    stats['code_blocks'] += len(inline_codes) * 2

    

    # 4. 表格: 含 | 的连续行

    table_lines = _re.findall(r'^\s*\|.+\|\s*$', text, _re.MULTILINE)

    if table_lines:

        stats['tables'] += len(table_lines) * 15  # 每行

        # 分隔行额外 +5

        for line in table_lines:

            if _re.match(r'^\s*\|[\s\-:|]+\|\s*$', line):

                stats['tables'] += 5

    

    # 5. 标题: # ~ ######

    headings = _re.findall(r'^#{1,6}\s+.+$', text, _re.MULTILINE)

    stats['headings'] = len(headings) * 3

    

    # 6. 列表项: - * + 开头，或 1. 2. 编号

    list_items = _re.findall(r'^(\s*)([-*+]|\d+\.)\s+', text, _re.MULTILINE)

    stats['lists'] = len(list_items) * 4

    

    # 7. 引用块: > 开头

    quotes = _re.findall(r'^>\s?.+$', text, _re.MULTILINE)

    stats['quotes'] = len(quotes) * 5

    

    # 8. 水平线: --- 或 *** 或 ___

    hrs = _re.findall(r'^(\s*)(---+|\*\*\*+|___+)(\s*)$', text, _re.MULTILINE)

    stats['hr'] = len(hrs) * 3

    

    # 9. 加粗: **text** 或 __text__

    bolds = _re.findall(r'\*\*[^*]+\*\*|__[^_]+__', text)

    stats['emphasis'] += len(bolds) * 3

    

    # 10. 斜体: *text* 或 _text_（排除已匹配的加粗）

    italics = _re.findall(r'(?<!\*)\*(?!\*)([^*]+?)(?<!\*)\*(?!\*)|(?<!_)_(?!_)([^_]+?)(?<!_)_(?!_)', text)

    stats['emphasis'] += len(italics) * 2

    

    # 11. 删除线: ~~text~~

    strikethroughs = _re.findall(r'~~[^~]+~~', text)

    stats['strikethrough'] = len(strikethroughs) * 3

    

    # 12. 脚注: [^n]

    footnotes = _re.findall(r'\[\^[^\]]+\]', text)

    stats['footnotes'] = len(footnotes) * 10

    

    # 13. 数学公式（块级）: $$...$$

    math_blocks = _re.findall(r'\$\$(.*?)\$\$', text, _re.DOTALL)

    for block in math_blocks:

        stats['math'] += 20 + block.count('\n') * 5

    

    # 14. 数学公式（行内）: $...$（排除已匹配的块级）

    inline_math = _re.findall(r'(?<!\$)\$(?!\$)([^$]+?)(?<!\$)\$(?!\$)', text)

    stats['math'] += len(inline_math) * 10

    

    # 15. HTML 标签

    html_tags = _re.findall(r'<([a-zA-Z][a-zA-Z0-9]*)[^>]*>', text)

    stats['html'] = len(html_tags) * 5

    

    # 16. Checkbox: - [ ] 或 - [x]

    checkboxes = _re.findall(r'-\s*\[[ xX]\]', text)

    stats['checkboxes'] = len(checkboxes) * 5

    

    # 17. 目录标记: [[toc]] 或 [TOC]

    tocs = _re.findall(r'\[\[toc\]\]|\[TOC\]', text, _re.IGNORECASE)

    stats['toc'] = len(tocs) * 8

    

    # --- 从原文本中移除 Markdown 标记，得到纯文本用于基础字符统计 ---

    clean = text

    # 移除代码块

    clean = _re.sub(r'```[^\n]*\n.*?```', '', clean, flags=_re.DOTALL)

    # 移除图片标记（保留 alt 文本作为纯文本，已在图片 token 中计费）

    clean = _re.sub(r'!\[([^\]]*)\]\([^)]+\)', r'\1', clean)

    # 移除链接（保留文本）

    clean = _re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', clean)

    # 移除表格分隔行

    clean = _re.sub(r'^\s*\|[\s\-:|]+\|\s*$', '', clean, flags=_re.MULTILINE)

    # 移除表格管道符

    clean = _re.sub(r'\|', ' ', clean)

    # 移除标题标记

    clean = _re.sub(r'^#{1,6}\s+', '', clean, flags=_re.MULTILINE)

    # 移除列表标记

    clean = _re.sub(r'^(\s*)([-*+]|\d+\.)\s+', r'\1', clean, flags=_re.MULTILINE)

    # 移除引用标记

    clean = _re.sub(r'^>\s?', '', clean, flags=_re.MULTILINE)

    # 移除水平线

    clean = _re.sub(r'^(\s*)(---+|\*\*\*+|___+)(\s*)$', '', clean, flags=_re.MULTILINE)

    # 移除格式标记（加粗/斜体/删除线）

    clean = _re.sub(r'\*\*([^*]+)\*\*', r'\1', clean)

    clean = _re.sub(r'__([^_]+)__', r'\1', clean)

    clean = _re.sub(r'\*([^*]+)\*', r'\1', clean)

    clean = _re.sub(r'_([^_]+)_', r'\1', clean)

    clean = _re.sub(r'~~([^~]+)~~', r'\1', clean)

    # 移除行内代码

    clean = _re.sub(r'`([^`]+)`', r'\1', clean)

    # 移除脚注

    clean = _re.sub(r'\[\^[^\]]+\]', '', clean)

    # 移除数学公式

    clean = _re.sub(r'\$\$(.*?)\$\$', '', clean, flags=_re.DOTALL)

    clean = _re.sub(r'(?<!\$)\$(?!\$)([^$]+?)(?<!\$)\$(?!\$)', '', clean)

    # 移除 HTML 标签

    clean = _re.sub(r'<[^>]+>', '', clean)

    # 移除 checkbox 标记

    clean = _re.sub(r'-\s*\[[ xX]\]', '-', clean)

    # 移除目录标记

    clean = _re.sub(r'\[\[toc\]\]|\[TOC\]', '', clean, flags=_re.IGNORECASE)

    

    # --- 基础字符统计 ---

    # 中文

    chinese_chars = _re.findall(r'[\u4e00-\u9fff\u3400-\u4dbf\U00020000-\U0002a6df]', clean)

    stats['chinese'] = len(chinese_chars) * 1.5

    

    # 英文单词

    english_words = _re.findall(r'[a-zA-Z]+(?:\'[a-zA-Z]+)?', clean)

    stats['english'] = len(english_words) * 1.0

    

    # 数字串

    number_groups = _re.findall(r'\d[\d,._]*\d|\d', clean)

    stats['numbers'] = len(number_groups) * 0.5

    

    # 标点

    punct = _re.findall(r'[，。！？、；：\u201c\u201d\u2018\u2019（）【】《》…—·\u3000\uff01\uff0c\uff0e\uff1f\uff1b\uff1a\uff08\uff09\u300a\u300b\u2026\u2014\u00b7,.\-!?;:\'"()\[\]{}]', clean)

    stats['punctuation'] = len(punct) * 0.5

    

    # 空白

    whitespace = _re.findall(r'\s', clean)

    stats['whitespace'] = len(whitespace) * 0.3

    

    # 汇总

    stats['total'] = sum(v for k, v in stats.items() if k != 'total')

    # 四舍五入到整数

    for k in stats:

        stats[k] = int(round(stats[k]))

    

    return stats





def _empty_token_result():

    return {

        'total': 0, 'chinese': 0, 'english': 0, 'numbers': 0,

        'punctuation': 0, 'whitespace': 0, 'tables': 0, 'code_blocks': 0,

        'links': 0, 'images': 0, 'lists': 0, 'headings': 0,

        'emphasis': 0, 'quotes': 0, 'hr': 0, 'footnotes': 0,

        'math': 0, 'html': 0, 'strikethrough': 0,

        'checkboxes': 0, 'toc': 0,

    }





def format_token_summary(stats, label=""):

    """格式化 token 估算结果为可读字符串"""

    if not stats or stats['total'] == 0:

        return "0 tokens"

    

    prefix = f"[{label}] " if label else ""

    lines = [

        f"{prefix}预估 Token 消耗: ~{stats['total']:,}",

    ]

    

    # Top 分类

    cats = []

    for name, key, unit in [

        ('中文', 'chinese', ''), ('英文', 'english', '词'),

        ('代码块', 'code_blocks', ''), ('表格', 'tables', ''),

        ('图片', 'images', ''), ('链接', 'links', ''),

        ('数学公式', 'math', ''), ('HTML标签', 'html', ''),

    ]:

        v = stats.get(key, 0)

        if v > 0:

            cats.append(f"{name}: {v:,}{unit}")

    

    if cats:

        lines.append("  " + " | ".join(cats[:4]))

        if len(cats) > 4:

            lines.append("  " + " | ".join(cats[4:]))

    

    return '\n'.join(lines)





# ========== 文档减肥器 ==========



class DocSlimmer:

    """文档减肥器 - 精简文本，去除冗余"""

    

    # 压缩词汇映射（先处理，避免与AI痕迹冲突）

    REDUNDANT_PATTERNS = [

        (r'\n{3,}', '\n\n'),

        (r' {2,}', ' '),

        (r'\t+', ' '),

        # 程度副词冗余

        (r'非常+', '很'),

        (r'特别+', '很'),

        (r'极其+', '很'),

        (r'十分+', '很'),

        (r'相当+', '很'),

        # 重复表达

        (r'真的', ''),

        (r'实际上', ''),

        (r'其实', ''),

        # 连接词冗余

        (r'也就是说', '即'),

        (r'并且', '且'),

        (r'同时', '且'),

        (r'但是', '但'),

        (r'然而', '但'),

        # 标点冗余

        (r'[。]{2,}', '。'),

        (r'[，]{2,}', '，'),

        (r'[！]{2,}', '！'),

        (r'[？]{2,}', '？'),

        # 清理连续标点+空格

        (r'[。！？，；]\s*', lambda m: m.group()[:-1] + ' ' if m.group().endswith('，') or m.group().endswith('。') else m.group()),

    ]

    

    # AI写作痕迹（后处理）

    AI_PATTERNS = [

        (r'首先、?其次、?再次、?最后', '第一、第二、第三'),

        (r'综上所述', '总之'),

        (r'值得注意的是', '但'),

        (r'可以说是不胜枚举', '不胜枚举'),

        (r'可以说', ''),

        (r'毋庸置疑', ''),

        (r'显而易见', ''),

        (r'不言而喻', ''),

        (r'众所周知', ''),

        (r'从某种意义上', ''),

        (r'总的来说', '总体'),

        (r'坦白地说', ''),

        (r'客观地说', ''),

        (r'想必', ''),

        (r'不由得', ''),

        (r'不禁', ''),

        (r'想必大家', ''),

        (r'那么我们', ''),

        (r'首先、其次、最后', '第一、第二'),

    ]

    

    def slim(self, text, compression_rate=0.3, remove_ai=False):

        if not text or not text.strip():

            return {"result": "", "stats": {}}

        

        original_length = len(text)

        result = text

        

        # 根据压缩强度决定应用的规则数量

        # compression_rate: 0.0-1.0，表示目标压缩比例

        # 0.0 = 不压缩，1.0 = 最大压缩

        

        # 基础清理（总是应用）

        for pattern, replacement in self.REDUNDANT_PATTERNS:

            if callable(replacement):

                result = re.sub(pattern, replacement, result)

            else:

                result = re.sub(pattern, replacement, result)

        

        # 根据压缩强度应用额外的压缩

        if compression_rate > 0.1:

            # 去除重复字符

            result = re.sub(r'(\w)\1{2,}', r'\1', result)

        

        if compression_rate > 0.3:

            # 移除多余的空格和制表符

            result = re.sub(r' {2,}', ' ', result)

            result = re.sub(r'\t+', ' ', result)

        

        if compression_rate > 0.5:

            # 移除某些冗余词汇

            result = re.sub(r'（[^）]*）', '', result)  # 移除括号内容

            result = re.sub(r'\([^)]*\)', '', result)   # 移除英文括号内容

        

        if compression_rate > 0.7:

            # 激进压缩：移除更多内容

            result = re.sub(r'【[^】]*】', '', result)  # 移除方括号

            result = re.sub(r'\[[^\]]*\]', '', result)

        

        # 去AI味

        if remove_ai:

            for pattern, replacement in self.AI_PATTERNS:

                result = re.sub(pattern, replacement, result)

        

        # 清理残留的连续标点和空格

        result = re.sub(r'[，。；！？\s]+', lambda m: ' ' if ' ' in m.group() and len(m.group()) > 1 else m.group(), result)

        result = result.strip()

        

        result = result.strip()

        new_length = len(result)

        actual_rate = (original_length - new_length) / original_length if original_length > 0 else 0

        

        return {

            "result": result,

            "stats": {

                "original_length": original_length,

                "new_length": new_length,

                "compression_rate": round(actual_rate * 100, 1),

                "reduced_chars": original_length - new_length

            }

        }





# ========== 文档脱敏器 ==========



class DocSanitizer:

    """文档脱敏器 - 去除个人信息，可选择性脱敏"""

    

    # 敏感信息正则（按分类）

    PATTERNS = {

        # ===== 个人敏感信息 =====

        '手机号': r'(?<!\d)1[3-9]\d[\s\-]?\d{4}[\s\-]?\d{4}(?!\d)',

        '邮箱': r'\b[\w.%+-]+@[\w.-]+\.[A-Za-z]{2,}\b',

        '身份证': r'(?<!\d)\d{17}[\dXx](?!\d)',

        '银行卡': r'(?<!\d)\d{16,19}(?!\d)',

        'IP地址': r'\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b',

        # ===== 商业敏感信息 =====

        '社会信用代码': r'(?<!\d)91\d{16}(?!\d)',        # 统一社会信用代码18位

        '营业执照号': r'(?<!\d)\d{15}(?!\d)',            # 15位工商注册号

        '开户许可证号': r'(?<!\d)[A-Z0-9]{14,16}(?!\d)',  # 14-16位开户许可

        '投标/成交价': r'[¥￥]?\d{1,12}(?:\.\d{2})?(?![元])',  # 金额（需配合上下文判断）

        '合同编号': r'(?<![a-zA-Z])[A-Z]{2,4}[-#]?\d{2,4}[-]?\d{2,8}(?![a-zA-Z])',  # HT/HT-2026-XXXX类

        '采购/订单编号': r'(?<![a-zA-Z])(?:CG|PO|DD|HT|FC)[-_]?\d{4,12}(?![a-zA-Z])',  # CG/PO/FC开头编号

        '固定电话': r'(?<!\d)0\d{2,3}[-]?\d{7,8}(?!\d)',  # 021-12345678

        '传真号': r'(?<!\d)(?:传真|Fax)[-:]?\s*0\d{2,3}[-]?\d{7,8}(?!\d)',

        '工号/学号': r'(?<![a-zA-Z])(?:工号|学号|员工号|编号)[-:]?\s*[A-Z0-9]{4,12}(?![a-zA-Z])',

        '项目代号': r'(?<![a-zA-Z])(?:项目[编号码]|PRJ|PROJ)[-_]?\d{2,8}(?![a-zA-Z])',

        '邮编': r'(?<!\d)\d{6}(?!\d)',

    }

    

    # 名称映射

    ITEM_LABELS = {

        '手机号': '手机号',

        '邮箱': '邮箱',

        '身份证': '身份证',

        '银行卡': '银行卡',

        'IP地址': 'IP地址',

    }

    

    @staticmethod

    def available_items():

        return list(DocSanitizer.PATTERNS.keys())

    

    def _mask_phone(self, m):

        p = m.group()

        return f"{p[:3]}****{p[7:]}"

    

    def _mask_email(self, m):

        e = m.group()

        parts = e.split('@')

        if len(parts) == 2:

            name, domain = parts

            n = len(name)

            if n >= 3:

                return f"{name[:2]}***@{domain}"

            elif n == 2:

                return f"{name[0]}***@{domain}"

            else:

                return f"***@{domain}"

        return e

    

    def _mask_id(self, m):

        i = m.group()

        return f"{i[:6]}**********{i[-1] if len(i)==18 else ''}"

    

    def _mask_bank(self, m):

        b = m.group()

        return f"{b[:4]}****{b[-4:]}"

    

    def _mask_ip(self, m):

        return 'xxx.xxx.xxx.xxx'

    

    def _mask_passport(self, m):

        p = m.group()

        if len(p) >= 5:

            return p[0] + '*' * (len(p) - 2) + p[-1]

        return '*' * len(p)

    

    def _mask_mac(self, m):

        mac = m.group()

        parts = mac.replace('-', ':').split(':')

        if len(parts) == 6:

            return ':'.join([parts[0]] + ['****'] + [parts[-1]])

        return '**:**:**:**:**:**'

    

    def _mask_imei(self, m):

        e = m.group()

        if len(e) >= 8:

            return e[:6] + '******' + e[-6:]

        return '*' * len(e)

    

    def _mask_plate(self, m):

        p = m.group()

        return p[:2] + '*' * (len(p) - 2)

    

    def _mask_social_card(self, m):

        c = m.group()

        return c[:4] + '*' * (len(c) - 8) + c[-4:]

    

    def _mask_medical_record(self, m):

        c = m.group()

        digits = re.sub(r'\D', '', c)

        prefix = re.sub(r'\d', '', c)

        return prefix + '*' * len(digits)

    

    def _mask_docnum(self, m):

        c = m.group()

        digits = re.sub(r'\D', '', c)

        prefix = re.sub(r'\d', '', c)

        return prefix + '*' * len(digits)

    

    def _mask_doclevel(self, m):

        c = m.group()

        if c.startswith('【') and c.endswith('】'):

            return '【' + '*' * (len(c) - 4) + '】'

        if '文件' in c:

            return '*' * len(c)

        return re.sub(r'[*☆★\d]+', '*', c)

    

    def _mask_docref(self, m):

        c = m.group()

        return re.sub(r'\d+(?=号)', lambda x: '*' * len(x.group()), c)

    

    def sanitize(self, text, custom_words=None, items=None):

        import sys
        print(f"[DEBUG] DocSanitizer.sanitize items={items!r}, len={len(text)}", file=sys.stderr)

        if not text:

            return {"result": "", "stats": {"total": 0}}

        

        stats = {}

        result = text

        

        if items is None:

            items = list(self.PATTERNS.keys())

        

        def mask_code(m):

            c = m.group()

            if len(c) <= 4:

                return '*' * len(c)

            half = len(c) // 2

            return f"{c[:half]}***{c[-half:]}"



        def mask_phone_last(m):

            p = m.group()

            digits = re.sub(r'\D', '', p)

            if len(digits) >= 8:

                return f"{digits[:3]}****{digits[-4:]}"

            return '****'



        def mask_social_credit(m):

            c = m.group()

            if len(c) == 18:

                return f"{c[:4]}**********{c[-4:]}"

            return mask_code(m)



        def mask_landline(m):

            p = m.group()

            digits = re.sub(r'\D', '', p)

            if len(digits) >= 8:

                return f"{digits[:3]}****{digits[-4:]}"

            return '****'



        def mask_fax(m):

            return '***-********'



        def mask_amount(m):

            return m.group()[:1] + '***'



        # ===== 个人敏感信息 =====

        if '手机号' in items:

            found = re.findall(self.PATTERNS['手机号'], result)

            stats['手机号'] = len(found)

            result = re.sub(self.PATTERNS['手机号'], self._mask_phone, result)



        if '邮箱' in items:

            found = re.findall(self.PATTERNS['邮箱'], result)

            stats['邮箱'] = len(found)

            result = re.sub(self.PATTERNS['邮箱'], self._mask_email, result)



        if '身份证' in items:

            found = re.findall(self.PATTERNS['身份证'], result)

            stats['身份证'] = len(found)

            result = re.sub(self.PATTERNS['身份证'], self._mask_id, result)



        if '银行卡' in items:

            found = re.findall(self.PATTERNS['银行卡'], result)

            stats['银行卡'] = len(found)

            result = re.sub(self.PATTERNS['银行卡'], self._mask_bank, result)



        if 'IP地址' in items:

            found = re.findall(self.PATTERNS['IP地址'], result)

            stats['IP地址'] = len(found)

            result = re.sub(self.PATTERNS['IP地址'], self._mask_ip, result)



        # ===== 商业敏感信息 =====

        if '社会信用代码' in items:

            found = re.findall(r'(?<![A-Z0-9])[0-9A-Z]{18}(?![A-Z0-9])', result)

            stats['社会信用代码'] = len(found)

            result = re.sub(r'(?<![A-Z0-9])[0-9A-Z]{18}(?![A-Z0-9])', mask_social_credit, result)



        if '营业执照号' in items:

            found = re.findall(r'(?<!\d)\d{15}(?!\d)', result)

            stats['营业执照号'] = len(found)

            result = re.sub(r'(?<!\d)\d{15}(?!\d)', lambda m: f"{m.group()[:5]}****{m.group()[-3:]}", result)



        if '开户许可证号' in items:

            found = re.findall(r'(?<![A-Z0-9])[A-Z0-9]{14,16}(?![A-Z0-9])', result)

            stats['开户许可证号'] = len(found)

            result = re.sub(r'(?<![A-Z0-9])[A-Z0-9]{14,16}(?![A-Z0-9])', mask_code, result)



        if '投标/成交价' in items:

            # 带货币符号的金额

            found_sym = re.findall(r'[¥￥$€£]\s*[\d,]+(?:\.\d+)?', result)

            stats['投标/成交价'] = len(found_sym)

            for m in found_sym:

                digits = re.sub(r'[^\d]', '', m)

                prefix = m[:len(m) - len(digits)]

                result = result.replace(m, prefix + '*' * min(len(digits), 8), 1)

            # USD格式

            found_usd = re.findall(r'USD\s*[\d,]+(?:\.\d+)?', result)

            for m in found_usd:

                digits = re.sub(r'[^\d]', '', m)

                result = result.replace(m, 'USD ' + '*' * min(len(digits), 6), 1)

            # 纯数字+单位

            found_unit = re.findall(r'[\d,]+(?:\.\d+)?\s*(?:万|亿|美元|欧元|英镑|元)', result)

            for m in found_unit:

                digits = re.sub(r'[^\d]', '', m)

                unit = re.sub(r'[\d,\s]', '', m)

                result = result.replace(m, '*' * min(len(digits), 8) + unit, 1)

            # 人民币中文大写金额

            found_cn = re.findall(r'人民币[零壹贰叁肆伍陆柒捌玖拾佰仟萬万〇\d整兆亿万]+元整?', result)

            for m in found_cn:

                result = result.replace(m, '人民币*元整', 1)

            if not found_sym and not found_usd and not found_unit and not found_cn:

                stats['投标/成交价'] = 0



        if '合同编号' in items:

            found = re.findall(r'[A-Z]{2,8}[-_/](?:[A-Z0-9]+[-_/])*[A-Z0-9]*\d{4,}[A-Z0-9]*(?:[-_/][A-Z0-9]+)*', result)

            stats['合同编号'] = len(found)

            result = re.sub(r'[A-Z]{2,8}[-_/](?:[A-Z0-9]+[-_/])*[A-Z0-9]*\d{4,}[A-Z0-9]*(?:[-_/][A-Z0-9]+)*', mask_code, result)



        if '采购/订单编号' in items:

            found = re.findall(r'(?<![A-Z])(?:CG|PO|DD|FC)[-_]?\d{4,12}(?![A-Z0-9])', result, re.IGNORECASE)

            stats['采购/订单编号'] = len(found)

            result = re.sub(r'(?<![A-Z])(?:CG|PO|DD|FC)[-_]?\d{4,12}(?![A-Z0-9])', mask_code, result, flags=re.IGNORECASE)



        if '固定电话' in items:

            found = re.findall(self.PATTERNS['固定电话'], result)

            stats['固定电话'] = len(found)

            result = re.sub(self.PATTERNS['固定电话'], mask_landline, result)



        if '传真号' in items:

            found = re.findall(self.PATTERNS['传真号'], result)

            stats['传真号'] = len(found)

            result = re.sub(self.PATTERNS['传真号'], mask_fax, result)



        if '工号/学号' in items:

            found = re.findall(self.PATTERNS['工号/学号'], result)

            stats['工号/学号'] = len(found)

            result = re.sub(self.PATTERNS['工号/学号'], mask_code, result)



        if '项目代号' in items:

            found = re.findall(r'(?<![a-zA-Z])(?:项目[编号码]|PRJ|PROJ)[-_]?\d{2,8}(?![a-zA-Z0-9])', result, re.IGNORECASE)

            stats['项目代号'] = len(found)

            result = re.sub(r'(?<![a-zA-Z])(?:项目[编号码]|PRJ|PROJ)[-_]?\d{2,8}(?![a-zA-Z0-9])', mask_code, result, flags=re.IGNORECASE)



        if '邮编' in items:

            found = re.findall(r'(?<!\d)\d{6}(?!\d)', result)

            stats['邮编'] = len(found)

            result = re.sub(r'(?<!\d)\d{6}(?!\d)', lambda m: f"{m.group()[:2]}****{m.group()[-1]}", result)



        # ===== 新增类型 =====

        if '护照号' in items:

            found = re.findall(r'\b[EeGgPpDdSsHhLl][A-Za-z]?\d{7,9}\b', result)

            stats['护照号'] = len(found)

            for m in found:

                result = result.replace(m, m[0] + '*' * (len(m) - 2) + m[-1], 1)



        if 'Mac地址' in items:

            found = re.findall(r'\b([0-9A-Fa-f]{2}[:-]){5}[0-9A-Fa-f]{2}\b', result)

            stats['Mac地址'] = len(found)

            for m in found:

                sep = ':' if ':' in m else '-'

                parts = m.split(sep)

                result = result.replace(m, sep.join([parts[0]] + ['****'] + [parts[-1]]), 1)



        if 'IMEI' in items:

            found = re.findall(r'(?<!\d)\d{15}(?!\d)', result)

            stats['IMEI'] = len(found)

            for m in found:

                result = result.replace(m, m[:6] + '******' + m[-6:], 1)



        if '车牌号' in items:

            found = re.findall(r'[京津沪渝冀豫云辽黑湘皖鲁新苏浙赣鄂桂甘晋蒙陕吉闽贵粤青藏川宁琼使领警][A-Z][·]?[A-HJ-NP-Z0-9]{4,5}[A-HJ-NP-Z0-9挂学港澳]?', result)

            stats['车牌号'] = len(found)

            for m in found:

                result = result.replace(m, m[:2] + '*' * (len(m) - 2), 1)



        if '社保卡号' in items:

            found = re.findall(r'(?<!\d)\d{14,18}(?!\d)', result)

            stats['社保卡号'] = len(found)

            for m in found:

                result = result.replace(m, m[:4] + '*' * (len(m) - 8) + m[-4:], 1)



        if '医保卡号' in items:

            found = re.findall(r'(?<!\d)\d{15,18}(?!\d)', result)

            stats['医保卡号'] = len(found)

            for m in found:

                result = result.replace(m, m[:4] + '*' * (len(m) - 6) + m[-4:], 1)



        if '病历号' in items:

            found = re.findall(r'(?:(?:BL|MR|EMR|MZ|ZY|门诊号?|住院号?|病历号?|病案号?)[：:\-]?\d{4,}|\b[A-Z]{2,3}[-]?\d{6,12}\b)', result)

            stats['病历号'] = len(found)

            for m in found:

                digits = re.sub(r'\D', '', m)

                prefix = re.sub(r'\d', '', m)

                result = result.replace(m, prefix + '*' * len(digits), 1)



        if '公文份号' in items:

            found = re.findall(r'(?:(?:No|NO|Nr|№)[-:\s]*\d{4,}|份号[：:\s]+\d+|第\d{2,4}[-]\d{4,}号|文件编号[：:\s]*[A-Z]{2,3}[-_]\d{2,4}[-_]\d{3,})', result)

            stats['公文份号'] = len(found)

            for m in found:

                digits = re.sub(r'\D', '', m)

                prefix = re.sub(r'\d', '', m)

                result = result.replace(m, prefix + '*' * len(digits), 1)



        if '公文密级' in items:

            found = re.findall(r'(?:【[绝密机密秘密内部]+】|绝密\s*[★☆]?\s*\d*\s*年?|[机密秘密内部]+(?:文件|资料|通知|信息))', result)

            stats['公文密级'] = len(found)

            for m in found:

                if m.startswith('【') and m.endswith('】'):

                    result = result.replace(m, '【' + '*' * (len(m) - 4) + '】', 1)

                elif '文件' in m or '资料' in m or '通知' in m or '信息' in m:

                    result = result.replace(m, '*' * len(m), 1)

                else:

                    result = result.replace(m, '*' * len(m), 1)



        if '公文文号' in items:

            found = re.findall(r'[\u4e00-\u9fa5]+发〔\d{4}〕\s*(?:第\s*)?\d+(?:号)?|[\u4e00-\u9fa5]+政办发〔\d{4}〕\s*(?:第\s*)?\d+(?:号)?', result)

            stats['公文文号'] = len(found)

            for m in found:

                result = result.replace(m, re.sub(r'\d+(?=号)', lambda x: '*' * len(x.group()), m), 1)



        # 自定义敏感词（始终处理）

        custom_count = 0

        if custom_words:

            for word in custom_words:

                if word.strip() and len(word) >= 2:

                    count = result.count(word)

                    if count > 0:

                        masked = word[0] + '*' * (len(word) - 2) + word[-1] if len(word) > 2 else '**'

                        result = result.replace(word, masked)

                        custom_count += count

        if custom_count > 0:

            stats['自定义'] = custom_count



        stats['总计'] = sum(stats.values())



        return {

            "result": result.strip(),

            "stats": stats

        }





# ========== 批量处理入口 ==========



def cmd_batch(args):

    """批量处理命令"""

    from batch_processor import batch_process, format_size

    import shutil



    # 确定输出目录

    if args.out_dir:

        out_dir = args.out_dir

    else:

        cfg = load_config()

        default_out = cfg.get('output_dir', 'output')

        if default_out == 'output':

            # 默认输出到源目录的父目录

            src_parent = str(Path(args.folder).resolve().parent)

            ts = datetime.now().strftime('%Y%m%d_%H%M%S')

            out_dir = os.path.join(src_parent, f"{Path(args.folder).name}_处理结果_{ts}")

        else:

            out_dir = default_out



    # 选项

    options = {}

    if args.command == 'batch-slim':

        options['compression'] = args.compression

        options['remove_ai'] = args.ai

        action = 'slim'

    else:

        options['custom_words'] = args.words or []

        options['sanitize_items'] = args.items or None

        action = 'sanitize'



    print(f"[批量{action}] 文件夹: {args.folder}")

    print(f"[输出] {out_dir}")

    print()



    result = batch_process(

        folder_path=args.folder,

        action=action,

        options=options,

        recursive=not args.no_recursive,

        workers=args.workers or 4,

        out_base=out_dir

    )



    if 'error' in result and result.get('total', 0) == 0:

        print(f"\n错误: {result['error']}")

        return



    s = result['summary']

    print()

    print(f"{'='*50}")

    print(f"  完成！耗时 {s['elapsed_seconds']}秒")

    print(f"  成功: {s['success']}  失败: {s['error']}  跳过: {s['skip']}")

    if s['total_input_size'] > 0:

        rate = s['size_reduced'] / s['total_input_size'] * 100

        print(f"  原始: {format_size(s['total_input_size'])} → 处理后: {format_size(s['total_output_size'])}")

        print(f"  节省: {format_size(s['size_reduced'])} ({rate:.1f}%)")

    print(f"  输出: {result['output_dir']}")

    print(f"  报告: {result['output_dir']}\\处理报告.txt")

    print(f"{'='*50}")



    # 打印报告

    print("\n" + result['report'])





# ========== 主程序 ==========



def load_config():

    """加载配置"""

    config_path = Path(__file__).parent / 'config.json'

    if config_path.exists():

        try:

            with open(config_path, 'r', encoding='utf-8') as f:

                return json.load(f)

        except Exception:

            pass

    return {}





def resolve_output_path(input_path, action, user_output, out_dir, fmt):

    """解析输出路径，支持：手动指定 / 默认output目录 / 同目录"""

    Path_like = type(Path())

    

    # 如果用户直接指定了完整路径

    if user_output:

        return Path(user_output)

    

    # 确定输出目录

    if out_dir:

        out = Path(out_dir)

    else:

        cfg = load_config()

        default_out = cfg.get('output_dir')

        if default_out:

            out = Path(__file__).parent / default_out

        else:

            # 默认：源文件同目录

            out = Path(input_path).parent

    

    # 创建输出目录（如果不存在）

    if not out.exists():

        try:

            out.mkdir(parents=True, exist_ok=True)

        except PermissionError:

            print(f"[提示] 无法创建输出目录 '{out}'，将输出到控制台", file=sys.stderr)

            out = None

        except Exception:

            out = None

    

    # 生成文件名

    src = Path(input_path)

    action_tag = '减肥' if action == 'slim' else '脱敏'

    

    if src.suffix.lower() in ['.docx', '.xlsx', '.xls', '.pptx', '.pdf']:

        # 格式文件：提取文本后保存为.ssd

        stem = src.stem

        out_name = f"{stem}_{action_tag}.ssd"

    else:

        # 文本文件：保持原扩展名

        stem = src.stem

        ext = fmt if fmt else src.suffix

        out_name = f"{stem}_{action_tag}{ext}"

    

    return out / out_name





def main():

    # 加载配置

    cfg = load_config()

    

    parser = argparse.ArgumentParser(

        description='SafeShrink v2.0 - 文档减肥 & 脱敏（支持Office/PDF）',

        formatter_class=argparse.RawDescriptionHelpFormatter,

        epilog='''

支持格式:

  文本: .txt .md .json .csv .xml .html .log

  Word: .docx  Excel: .xlsx .xls  PPT: .pptx  PDF: .pdf



保存策略:

  - 不覆盖源文件，自动另存

  - 默认保存到 output/ 子目录（可在 config.json 中修改）

  - 也可用 --out-dir 指定输出目录



示例:

  # 文档减肥（自动保存到 output/）

  safeshrink slim input.pdf

  safeshrink slim input.docx -o custom.txt

  safeshrink slim "要精简的文本" -o custom.txt

  

  # 文档脱敏（自动保存到 output/）

  safeshrink sanitize input.xlsx --out-dir D:\\输出

  safeshrink sanitize "文本..." --words 张三

  

  # 查看依赖

  safeshrink --check

        '''

    )

    

    parser.add_argument('--check', action='store_true', help='检查依赖状态')

    parser.add_argument('--version', action='version', version='SafeShrink v0.9.0')

    

    subparsers = parser.add_subparsers(dest='command', help='子命令')

    

    # 全局 JSON 输出标志（顶级参数，对 slim/sanitize 生效）

    parser.add_argument('--json', dest='json_output', action='store_true',

                        help='以 JSON 格式输出结果（供程序调用）')



    # slim

    sp = subparsers.add_parser('slim', help='文档减肥')

    sp.add_argument('text', nargs='?', help='要处理的文本')

    sp.add_argument('-i', '--input', help='输入文件')

    sp.add_argument('-o', '--output', help='输出文件路径（完整路径）')

    sp.add_argument('-d', '--out-dir', help='输出目录（默认使用 config.json 中的 output_dir）')

    sp.add_argument('-c', '--compression', type=float, default=0.3, help='压缩率 (0.0-1.0)')

    sp.add_argument('--ai', action='store_true', help='去除AI写作痕迹')

    sp.add_argument('--sheet', type=int, default=0, help='Excel sheet索引')

    sp.add_argument('--json', dest='json_output', action='store_true',

                    help='以 JSON 格式输出结果（供程序调用）')

    

    # sanitize

    ep = subparsers.add_parser('sanitize', help='文档脱敏')

    ep.add_argument('text', nargs='?', help='要处理的文本')

    ep.add_argument('-i', '--input', help='输入文件')

    ep.add_argument('-o', '--output', help='输出文件路径（完整路径）')

    ep.add_argument('-d', '--out-dir', help='输出目录')

    ep.add_argument('--words', nargs='*', help='自定义敏感词')

    ep.add_argument('--items', nargs='*', 

        help='指定要处理的脱敏项，可用值: 手机号,邮箱,身份证,银行卡,IP地址')

    ep.add_argument('--sheet', type=int, default=0, help='Excel sheet索引')

    ep.add_argument('--format', help='指定输出格式')

    ep.add_argument('--json', dest='json_output', action='store_true',

                    help='以 JSON 格式输出结果（供程序调用）')

    

    # batch-slim

    bp = subparsers.add_parser('batch-slim', help='批量文档减肥（文件夹）')

    bp.add_argument('folder', help='输入文件夹路径')

    bp.add_argument('-o', '--out-dir', help='输出目录')

    bp.add_argument('-c', '--compression', type=float, default=0.3, help='压缩率')

    bp.add_argument('--ai', action='store_true', help='去除AI味')

    bp.add_argument('-w', '--workers', type=int, help='并行线程数（默认4）')

    bp.add_argument('--no-recursive', action='store_true', help='不递归子文件夹')

    

    # batch-sanitize

    bx = subparsers.add_parser('batch-sanitize', help='批量文档脱敏（文件夹）')

    bx.add_argument('folder', help='输入文件夹路径')

    bx.add_argument('-o', '--out-dir', help='输出目录')

    bx.add_argument('--words', nargs='*', help='自定义敏感词')

    bx.add_argument('--items', nargs='*', help='指定脱敏项')

    bx.add_argument('-w', '--workers', type=int, help='并行线程数（默认4）')

    bx.add_argument('--no-recursive', action='store_true', help='不递归子文件夹')

    

    args = parser.parse_args()

    

    # 检查依赖

    if args.check:

        print("=== 依赖检查 ===")

        deps_list = [

            ('docx', 'python-docx', 'Word .docx'),

            ('openpyxl', 'openpyxl', 'Excel .xlsx'),

            ('pptx', 'python-pptx', 'PPT .pptx'),

            ('pdfplumber', 'pdfplumber', 'PDF (表格)'),

            ('pymupdf', 'PyMuPDF', 'PDF (通用)'),

            ('xlrd', 'xlrd', 'Excel .xls'),

        ]

        for key, pkg, desc in deps_list:

            status = "[OK]" if DEPS.get(key) else "[MISSING]"

            print(f"  {status} {desc:12} ({pkg})")

        print(f"\n[配置] output_dir: {cfg.get('output_dir', '未设置（输出到源文件目录）')}")

        return

    

    if not args.command:

        parser.print_help()

        return

    

    # 批量处理

    if args.command in ('batch-slim', 'batch-sanitize'):

        import shutil

        from batch_processor import batch_process, format_size



        if args.out_dir:

            out_dir = args.out_dir

        else:

            src_parent = str(Path(args.folder).resolve().parent)

            ts = datetime.now().strftime('%Y%m%d_%H%M%S')

            out_dir = os.path.join(src_parent, f"{Path(args.folder).name}_处理结果_{ts}")



        options = {}

        if args.command == 'batch-slim':

            options['compression'] = args.compression

            options['remove_ai'] = args.ai

            options['output_ext'] = '.ssd'

            action = 'slim'

        else:

            options['custom_words'] = args.words or []

            options['sanitize_items'] = args.items or None

            options['output_ext'] = '.ssd'

            action = 'sanitize'



        print(f"[批量{action}] 文件夹: {args.folder}")

        print(f"[输出] {out_dir}\n")



        result = batch_process(

            folder_path=args.folder,

            action=action,

            options=options,

            recursive=not args.no_recursive,

            workers=args.workers or 4,

            out_base=out_dir

        )



        if 'error' in result and result.get('total', 0) == 0:

            print(f"\n错误: {result['error']}")

            return



        s = result['summary']

        print(f"\n{'='*50}")

        print(f"  完成！耗时 {s['elapsed_seconds']}秒")

        print(f"  成功: {s['success']}  失败: {s['error']}  跳过: {s['skip']}")

        if s['total_input_size'] > 0:

            rate = s['size_reduced'] / s['total_input_size'] * 100

            print(f"  原始: {format_size(s['total_input_size'])} → 处理后: {format_size(s['total_output_size'])}")

            print(f"  节省: {format_size(s['size_reduced'])} ({rate:.1f}%)")

        print(f"  输出: {result['output_dir']}")

        print(f"  报告: {result['output_dir']}\\处理报告.txt")

        print(f"{'='*50}\n{result['report']}")

        return

    

    # 获取输入文本

    text = None

    input_file = None

    if args.input:

        input_file = args.input

        opts = {'sheet': args.sheet}

        try:

            text = read_file(args.input, opts)

            print(f"[读入] {args.input} ({len(text)} 字符)")

        except ImportError as e:

            print(f"[错误] {e}", file=sys.stderr)

            sys.exit(1)

        except Exception as e:

            print(f"[错误] 读取失败: {e}", file=sys.stderr)

            sys.exit(1)

    elif args.text:

        text = args.text

    else:

        print("[错误] 请提供文本或输入文件 (-i)", file=sys.stderr)

        sys.exit(1)

    

    # 处理

    json_mode = getattr(args, 'json_output', False)

    stats = {}



    if args.command == 'slim':

        processor = DocSlimmer()

        result = processor.slim(text, args.compression, args.ai)

        stats = result['stats']

        if not json_mode:

            print(f"[减肥] 压缩率: {stats.get('compression_rate', 0)}%")

            print(f"       减少字符: {stats.get('reduced_chars', 0)}")

        

    elif args.command == 'sanitize':

        processor = DocSanitizer()

        sanitize_items = getattr(args, 'items', None)

        result = processor.sanitize(text, args.words, sanitize_items)

        stats = result['stats']

        if not json_mode:

            total = stats.get('总计', 0)

            detail = ', '.join(f"{k}:{v}" for k, v in stats.items() if k != '总计' and v > 0)

            print(f"[脱敏] 共脱敏 {total} 项: {detail if detail else '无'}")

    

    # 输出（核心：永远不覆盖源文件）

    output_path = resolve_output_path(

        input_path=input_file or '.',

        action=args.command,

        user_output=args.output,

        out_dir=args.out_dir,

        fmt=args.format if hasattr(args, 'format') else None

    )

    

    # 如果输出目录无法创建

    if output_path is None:

        if json_mode:

            print(json.dumps({"success": True, "output_path": None, "stats": stats,

                              "result": result['result'][:10000]}, ensure_ascii=False))

        else:

            print('\n' + '='*50)

            print(result['result'][:3000])

            if len(result['result']) > 3000:

                print(f"\n... (共 {len(result['result'])} 字符，已截断)")

        return

    

    # 写入文件

    write_ok = True

    write_error = None

    actual_output = str(output_path)

    try:

        fmt = args.format if (hasattr(args, 'format') and args.format) else output_path.suffix.lstrip('.') or 'txt'

        write_file(str(output_path), result['result'], fmt)

        if not json_mode:

            print(f"[保存] -> {output_path}")

    except ImportError as e:

        write_error = str(e)

        txt_path = str(output_path) + '.txt'

        write_txt(txt_path, result['result'])

        actual_output = txt_path

        if not json_mode:

            print(f"[保存] -> {txt_path} (格式降级)")

    except Exception as e:

        write_ok = False

        write_error = str(e)

        actual_output = None

        if not json_mode:

            print(f"[错误] 保存失败: {e}", file=sys.stderr)

            print("[提示] 权限不足时，可尝试指定其他输出目录：--out-dir D:\\输出")

            print('\n' + '='*50)

            print(result['result'][:3000])

    

    # JSON 输出

    if json_mode:

        json_out = {

            "success": write_ok,

            "action": args.command,

            "output_path": actual_output,

            "stats": stats,

        }

        if input_file:

            json_out["input_file"] = input_file

        print(json.dumps(json_out, ensure_ascii=False))





# ========== 图片压缩 ==========



def compress_image(input_path, output_path=None, quality=85, max_size=None):

    """

    压缩图片文件

    

    Args:

        input_path: 输入图片路径

        output_path: 输出图片路径，None 则覆盖原文件

        quality: JPEG 质量 (1-100)

        max_size: 最大尺寸 (width, height)，None 则保持原尺寸

    

    Returns:

        dict: {'success': bool, 'original_size': int, 'new_size': int, 'output_path': str}

    """

    from PIL import Image

    import os

    

    try:

        # 获取原文件大小

        original_size = os.path.getsize(input_path)

        

        # 打开图片

        img = Image.open(input_path)

        

        # 转换 RGBA/RGB

        if img.mode in ('RGBA', 'P'):

            # 保持 RGBA 用于 PNG

            pass

        elif img.mode != 'RGB':

            img = img.convert('RGB')

        

        # 调整尺寸

        if max_size:

            img.thumbnail(max_size, Image.Resampling.LANCZOS)

        

        # 确定输出路径

        if output_path is None:

            output_path = input_path

        

        # 确定格式

        input_ext = os.path.splitext(input_path)[1].lower()

        

        if input_ext in ['.jpg', '.jpeg']:

            img.save(output_path, 'JPEG', quality=quality, optimize=True)

        elif input_ext == '.png':

            # PNG 使用不同的压缩方式

            img.save(output_path, 'PNG', optimize=True)

        elif input_ext == '.gif':

            img.save(output_path, 'GIF', optimize=True)

        else:

            img.save(output_path, quality=quality, optimize=True)

        

        # 获取新文件大小

        new_size = os.path.getsize(output_path)

        

        return {

            'success': True,

            'original_size': original_size,

            'new_size': new_size,

            'output_path': output_path,

            'saved': original_size - new_size,

            'saved_percent': round((1 - new_size/original_size) * 100, 1) if original_size > 0 else 0

        }

        

    except Exception as e:

        return {

            'success': False,

            'error': str(e)

        }





def get_image_info(path):

    """获取图片信息"""

    from PIL import Image

    import os

    

    try:

        img = Image.open(path)

        return {

            'width': img.width,

            'height': img.height,

            'mode': img.mode,

            'format': img.format,

            'size': os.path.getsize(path),

            'size_str': format_size(os.path.getsize(path))

        }

    except Exception as e:

        return {'error': str(e)}





if __name__ == "__main__":

    main()

